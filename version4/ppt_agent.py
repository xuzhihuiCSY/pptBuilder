from __future__ import annotations
import os, json, re
from typing import TypedDict, Optional, List, Dict, Any, Literal

# --- LangChain / LangGraph ---
from transformers import AutoModelForCausalLM, AutoTokenizer, pipeline, AutoConfig
import torch
from langchain_huggingface import ChatHuggingFace, HuggingFacePipeline
from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from langchain_core.runnables import RunnableLambda
from langgraph.graph import StateGraph, END

# --- Tools ---
from pydantic import BaseModel, Field
from langchain.tools import StructuredTool
import fitz  # PyMuPDF
from pptx import Presentation
import pathlib
import math
from functools import lru_cache

GENERIC_TITLE_PAT = re.compile(r"^(topic|section|slide)\s*\d+\b", re.I)

try:
    from transformers import BitsAndBytesConfig
except Exception:
    BitsAndBytesConfig = None
# =============================================================
# 1) Tool Implementations (pure Python + small LLM helpers)
# =============================================================

class ReadPDFIn(BaseModel):
    path: str = Field(..., description="Local path to a PDF file")
    pages: Optional[str] = Field(default=None, description='Optional page ranges like "1-3,5"; omit to read all.')
    mode: Optional[str] = Field(default="fast", description='Reading mode: "fast" or "full".')
    verbose: Optional[bool] = Field(default=True, description="Print progress while reading")

@lru_cache(maxsize=1)
def get_llm_singleton(temp: float = 0.2, num_predict: int = 320):
    return get_llm_20b(temp=temp, num_predict=num_predict)

def get_llm_20b(temp: float = 0.2, num_predict: int = 320):
    import os, warnings
    model_id = os.getenv("HF_MODEL_ID", "Qwen/Qwen2.5-7B-Instruct")
    want_4bit = os.getenv("HF_LOAD_4BIT", "0") == "1"

    tok = AutoTokenizer.from_pretrained(model_id, use_fast=True)
    cfg = AutoConfig.from_pretrained(model_id)

    qcfg = getattr(cfg, "quantization_config", None)
    is_mxfp4 = False
    if qcfg:
        qt = getattr(qcfg, "quant_method", None) or getattr(qcfg, "quantization_method", None) or getattr(qcfg, "quant_type", None)
        is_mxfp4 = str(qt).lower().startswith("mxfp4")

    load_kwargs = dict(device_map="auto")
    if is_mxfp4:
        load_kwargs.update(dtype=torch.bfloat16)
        model = AutoModelForCausalLM.from_pretrained(model_id, **load_kwargs)
    else:
        if want_4bit and BitsAndBytesConfig is not None:
            quant_cfg = BitsAndBytesConfig(
                load_in_4bit=True,
                bnb_4bit_use_double_quant=True,
                bnb_4bit_quant_type="nf4",
                bnb_4bit_compute_dtype=torch.bfloat16 if torch.cuda.is_available() else torch.float32,
            )
            model = AutoModelForCausalLM.from_pretrained(
                model_id,
                quantization_config=quant_cfg,
                device_map="auto",
            )
        else:
            dtype = torch.float16 if torch.cuda.is_available() else torch.float32
            model = AutoModelForCausalLM.from_pretrained(model_id, dtype=dtype, device_map="auto")

    gen_pipe = pipeline(
        task="text-generation",
        model=model,
        tokenizer=tok,
        max_new_tokens=int(os.getenv("HF_MAX_NEW_TOKENS", str(num_predict))),
        do_sample=(temp > 0),
        temperature=temp,
        top_p=0.9,
        repetition_penalty=1.05,
        pad_token_id=tok.eos_token_id,
        eos_token_id=tok.eos_token_id,
    )

    hf_llm = HuggingFacePipeline(pipeline=gen_pipe)
    return ChatHuggingFace(llm=hf_llm)


def read_pdf(path: str,
             pages: Optional[str] = None,
             mode: str = "full",
             max_chars: Optional[int] = None,
             max_pages: Optional[int] = None,
             verbose: bool = True) -> str:
    import os, fitz, re
    if verbose:
        print(f"[read_pdf] path: {path}")
    if not os.path.exists(path):
        raise FileNotFoundError(f"PDF not found: {path}")

    doc = fitz.open(path)

    def parse_ranges(spec: str) -> list[int]:
        out = set()
        for part in (spec or "").split(","):
            part = part.strip()
            if not part:
                continue
            if "-" in part:
                a, b = part.split("-", 1)
                a = max(1, int(a)); b = max(1, int(b))
                if a > b: a, b = b, a
                out.update(range(a - 1, min(b, len(doc))))
            else:
                p = max(1, int(part)) - 1
                if 0 <= p < len(doc):
                    out.add(p)
        return sorted(out)

    # === Unlimited by default ===
    if pages:
        selected = parse_ranges(pages)
    else:
        selected = list(range(len(doc)))  # ALL pages

    total_pages = len(selected)
    if verbose:
        cap_note = []
        if max_chars is not None: cap_note.append(f"cap≈{max_chars} chars")
        if max_pages is not None: cap_note.append(f"cap≈{max_pages} pages")
        note = f" ({', '.join(cap_note)})" if cap_note else ""
        print(f"[read_pdf] reading {total_pages} page(s) (mode={mode}){note}")

    chunks, total = [], 0
    for idx, pno in enumerate(selected, 1):
        page = doc.load_page(pno)
        txt = page.get_text("text") or ""
        if not txt.strip():
            blocks = page.get_text("blocks")
            if isinstance(blocks, list):
                txt = "\n".join(b[4] for b in blocks if len(b) >= 5)

        if txt:
            chunks.append(txt)
            total += len(txt)

        if verbose:
            print(f"[read_pdf] pages: {idx}/{total_pages} … total_chars={total}")

        # Early stops only if caps are provided
        if (max_pages is not None and idx >= max_pages) or (max_chars is not None and total >= max_chars):
            if verbose:
                print("[read_pdf] early stop (cap reached)")
            break

    return "\n".join(chunks)


ReadPDF = StructuredTool.from_function(
    func=read_pdf,
    name="read_pdf",
    description=(
        "Read text from a research PDF. Use when the user provided a PDF path and you need its content."
    ),
    args_schema=ReadPDFIn,
)

# --- Summarize (uses the same Ollama LLM behind the scenes) ---
class SummarizeIn(BaseModel):
    text: str = Field(..., description="Raw content (PDF text or notes)")
    style: Literal["slides", "prose", "exec"] = Field(default="slides", description="Output style")
    target_words: int = Field(400, description="Rough target length")
    verbose: Optional[bool] = Field(default=True, description="Print progress during summarization")


def summarize_text(text: str,
                   style: str = "summary",
                   target_words: int = 500,
                   verbose: bool = True) -> str:
    """Unlimited map→reduce summarizer; retries + strong fallback to avoid empty outputs."""
    import os, re
    SOURCE_CAP  = int(os.getenv("SUMMARY_SOURCE_CAP", "0"))     # 0 = no cap
    CHUNK_SIZE  = int(os.getenv("SUMMARY_CHUNK_SIZE", "3000"))
    MAX_CHUNKS  = int(os.getenv("SUMMARY_MAX_CHUNKS", "0"))     # 0 = unlimited

    src = (text or "")
    if SOURCE_CAP > 0:
        src = src[:SOURCE_CAP]
    if not src.strip():
        return ""

    chunks = _chunk_by_sentences(src, max_chars=CHUNK_SIZE)
    if MAX_CHUNKS > 0:
        chunks = chunks[:MAX_CHUNKS]
    if not chunks:
        return ""

    if verbose:
        cap_note = []
        if SOURCE_CAP > 0: cap_note.append(f"cap {SOURCE_CAP} chars")
        if MAX_CHUNKS > 0: cap_note.append(f"max {MAX_CHUNKS} chunks")
        note = f" ({', '.join(cap_note)})" if cap_note else ""
        print(f"[summ] map: {len(chunks)} chunk(s){note}")

    llm_map = get_llm_singleton(temp=0.2, num_predict=320)
    partials = []
    map_prompt_tpl = (
        "Summarize the following technical text in ~120 words. "
        "Write compact prose (no bullets, no headings). "
        "Keep problem, method, data, core results, and limitations. "
        "PRESERVE ALL NUMBERS (percentages, counts, dataset sizes, model names, metric values) verbatim.\n\n"
        "Text:\n{chunk}"
    )
    for i, ch in enumerate(chunks, 1):
        if verbose:
            print(f"[summ] map {i}/{len(chunks)} … (len={len(ch)})")
        resp = llm_map.invoke(map_prompt_tpl.format(chunk=ch)).content.strip()
        partials.append(resp)

    if verbose:
        print(f"[summ] reduce: combining {len(partials)} partial(s) …")

    def _reduce(num_predict: int, tgt_words: int, stronger: bool = False) -> str:
        llm_red = get_llm_singleton(temp=0.2, num_predict=num_predict)
        tgt = max(250, min(900, tgt_words))
        if stronger:
            reduce_prompt = (
                f"Combine the partial summaries into a single cohesive ~{tgt}-word executive summary.\n"
                "Write actual content (no prefaces), no bullets/headings, no meta commentary.\n"
                "CRITICAL: Preserve all quantitative details (percentages, counts, dataset sizes, parameters, metrics like AUROC) verbatim when mentioned.\n"
                "Include model/dataset names as written.\n\n" + "\n\n".join(partials)
            )
        else:
            reduce_prompt = (
                f"Combine the partial summaries into a single {tgt}-word executive summary.\n"
                "No bullets or headings. Avoid repetition.\n"
                "CRITICAL: Preserve percentages, counts, dataset sizes, model names, and metric values exactly.\n\n"
                + "\n\n".join(partials)
            )
        out = llm_red.invoke(reduce_prompt).content.strip()
        return re.sub(r"\s+", " ", out).strip()

    # Try reduce → retry bigger → fallback
    final = _reduce(num_predict=900, tgt_words=target_words, stronger=False)
    if len(final.split()) < 80:
        if verbose:
            print("[summ] reduce too short; retrying with larger output budget …")
        final_retry = _reduce(num_predict=1600, tgt_words=max(400, target_words), stronger=True)
        if len(final_retry.split()) > len(final.split()):
            final = final_retry

    # Fallback: join top partials and squeeze if still too short
    if len(final.split()) < 40:
        if verbose:
            print("[summ] reduce empty — using fallback join+squeeze …")
        joined = " ".join(partials[: min(10, len(partials))])
        # One last pass to compress joined partials into a single paragraph
        llm_fallback = get_llm_singleton(temp=0.2, num_predict=1200)
        fb_prompt = (
            "Condense the following text into a clear 300–450 word executive summary. "
            "No bullets or headings; write fluent paragraphs; avoid repetition:\n\n" + joined
        )
        final2 = llm_fallback.invoke(fb_prompt).content.strip()
        final2 = re.sub(r"\s+", " ", final2)
        if len(final2.split()) > len(final.split()):
            final = final2

    if verbose:
        print(f"[summ] done. ({len(final.split())} words)")

    facts = extract_quant_facts(src, max_items=12)
    if facts:
        final = final.rstrip() + "\n\nKey numbers: " + "; ".join(facts[:12])

    return final


Summarize = StructuredTool.from_function(
    func=summarize_text,
    name="summarize_text",
    description=(
        "Summarize content into concise bullets or prose. Use for 'summary only' requests or before making slides."
    ),
    args_schema=SummarizeIn,
)

# --- Slide outline JSON (use an LLM to structure nicely) ---
class OutlineIn(BaseModel):
    notes: Optional[str] = Field(default=None, description="If provided, use as source notes for the outline.")
    topic: Optional[str] = Field(default=None, description="If notes are not provided, generate an outline directly from this topic.")
    slide_target: int = Field(11, ge=5, le=30, description="Desired number of slides.")
    verbose: bool = Field(default=True, description="If true, print per-step outline progress to console.")

def densify_slide_with_llm(title: str, topic: str, notes: str, max_points: int = 5) -> tuple[str, list[str]]:
    llm = get_llm_singleton(temp=0.3, num_predict=400)
    prompt = (
        "Improve ONE presentation slide.\n"
        "Return STRICT JSON ONLY:\n"
        '{ "title": "Specific, informative title", "details": ["point 1", "point 2", "..."] }\n'
        "Rules:\n"
        "- 3 to 5 details; short phrases; no markdown bullets (*, -, •), no numbering.\n"
        "- Do NOT use phrases like 'Key takeaways include:' or 'Conclusion/Key Takeaways'.\n"
        "- Title must be informative (avoid 'Topic/Section/Slide N').\n\n"
        f"GlobalTopic: {topic}\n"
        f"SlideTitle: {title}\n\n"
        "Relevant notes:\n" + notes[:3500]
    )
    raw = llm.invoke(prompt).content.strip()
    m = re.search(r"\{[\s\S]*\}$", raw)
    js = m.group(0) if m else raw
    try:
        data = json.loads(js)
        t = re.sub(r"^[\s\*\-•·]+", "", str(data.get("title",""))).strip() or title
        pts = [re.sub(r"^[\s\*\-•·]+", "", str(x)).strip() for x in (data.get("details") or [])]
        pts = [p for p in pts if p][:max_points]
        if len(pts) < 3:
            raise ValueError("too few points")
        return t, pts
    except Exception:
        return title, synthesize_details_from_notes(title, notes, k=4)


def slide_outline_json(notes: Optional[str] = None,
                       slide_target: int = 11,
                       topic: Optional[str] = None,
                       verbose: bool = True) -> Dict[str, Any]:
    import json, os, re, datetime

    def log(msg: str):
        if verbose:
            print(msg)

    llm = get_llm_singleton(temp=0.2, num_predict=900)

    def _clean(s: str) -> str:
        return re.sub(r"^[\s\*\-•·]+", "", str(s or "")).strip()

    def extract_json(text: str) -> str:
        m = re.search(r"\{[\s\S]*\}$", text.strip())
        return (m.group(0) if m else text).strip()

    # ---------- Build prompt (two paths) ----------
    raw = ""  # always defined
    if notes and notes.strip():
        log(f"[outline] mode=notes  target={slide_target}")
        schema_text = (
            '{ "deck_title": "string", '
            '"slides": [ { "title": "string", "subtitle": "string or empty", '
            '"details": ["short phrase", "..."] } ] }'
        )
        prompt = (
            "You are an expert presentation writer.\n"
            "Return JSON ONLY matching this schema:\n" + schema_text + "\n\n"
            f"- EXACTLY {slide_target} slides.\n"
            "- Slide 1: Welcome/Title slide with a strong deck_title (also use it as slide 1 title) + short subtitle; details empty.\n"
            "- Slides 2..(N-2): 3–5 concise, factual details each.\n"
            "- Titles must be specific (avoid 'Topic 1/2/...'). Prefer titles hinting key findings or comparisons.\n"
            "- Slide (N-1): Conclusion (3–5 takeaways).\n"
            "- Slide N: Thanks for Watching (no details).\n"
            "- Avoid markdown bullets (*, -, •). JSON only, no extra text.\n"
            "- CRITICAL: When the notes include quantitative data (percentages, counts, dataset sizes, AUROC etc.), "
            "  preserve numbers verbatim and include at least one quantitative bullet per applicable content slide.\n\n"
            "NOTES:\n" + notes.strip()[:6000]
        )
        log("[outline] drafting raw JSON from notes …")
        try:
            raw = llm.invoke(prompt).content.strip()
        except Exception:
            raw = '{"deck_title":"Presentation","slides":[]}'
    else:
        topic = (topic or "").strip() or "General overview"
        log(f"[outline] mode=topic  target={slide_target}  topic={topic}")
        schema_text = (
            '{ "deck_title": "string", '
            '"slides": [ { "title": "string", "subtitle": "string or empty", '
            '"details": ["short phrase", "..."] } ] }'
        )
        prompt = (
            "You are a presentation planner. Return JSON ONLY matching this schema:\n" + schema_text + "\n\n"
            f"- EXACTLY {slide_target} slides.\n"
            "- Slide 1: Welcome/Title slide with a strong deck_title (use it as slide 1 title) + short subtitle; details empty.\n"
            "- Slides 2..(N-2): 3–5 concise, factual details each. Titles must be specific (no 'Topic 1/2/...').\n"
            "- Slide (N-1): Conclusion (3–5 takeaways).\n"
            "- Slide N: Thanks for Watching (no details).\n"
            "- Avoid markdown bullets (*, -, •). JSON only, no extra text.\n"
            "- Focus on the given TOPIC (do not produce generic presentation advice).\n\n"
            f"TOPIC: {topic}\n"
        )
        log("[outline] drafting raw JSON from topic …")
        try:
            raw = llm.invoke(prompt).content.strip()
        except Exception:
            raw = '{"deck_title":"Presentation","slides":[]}'

    # ---------- Parse JSON robustly ----------
    log("[outline] parsing JSON …")
    json_text = extract_json(raw)
    try:
        data = json.loads(json_text)
    except Exception:
        fixed = re.sub(r",\s*([}\]])", r"\1", json_text)
        try:
            data = json.loads(fixed)
        except Exception:
            data = {"deck_title": (topic or "Presentation"), "slides": []}

    deck_title = _clean(data.get("deck_title") or (topic or "Presentation"))
    slides = data.get("slides") or []
    log(f"[outline] parsed slides: {len(slides)}")

    # ---------- Normalize slides & enforce structure ----------
    cleaned: List[Dict[str, Any]] = []
    total = slide_target
    log(f"[outline] normalizing {total} slides …")
    for i in range(1, total + 1):
        # take from model output if present; else create a stub
        sl = slides[i-1] if i-1 < len(slides) else {}
        title = _clean(sl.get("title", f"Slide {i}"))
        subtitle = _clean(sl.get("subtitle", ""))
        details = [_clean(x) for x in (sl.get("details") or []) if _clean(x)]

        if i == 1:
            title = deck_title
            details = []
        elif i == max(2, total - 1):
            if not re.search(r"conclusion|takeaway|summary", title, re.I):
                title = "Conclusion"
            if len(details) < 3:
                details = ["(add point)"] * 3
        elif i == total:
            title = "Thanks for Watching"
            details = []
        else:
            if len(details) < 3:
                details += ["(add point)"] * (3 - len(details))
            details = details[:5]

        cleaned.append({"title": title, "subtitle": subtitle, "details": details})
        log(f"[outline] slide {i}/{total}: {title}")

    # ---------- Inject quantitative bullets if we have notes ----------
    if notes and notes.strip():
        log("[outline] injecting quantitative bullets …")
        quant_pool = extract_quant_facts(notes, max_items=40)
        if quant_pool:
            qi = 0
            for i in range(2, max(2, total - 1)):        # slides 2..(N-2)
                det = cleaned[i-1].get("details") or []
                has_number = any(re.search(r"\d", d) for d in det)
                if not has_number and qi < len(quant_pool):
                    det = det[:4] + [quant_pool[qi][:120]]
                    qi += 1
                cleaned[i-1]["details"] = det
            log(f"[outline] injected numeric bullets: {qi}")

    # ---------- Polish ----------
    src_for_polish = (notes or topic or "")
    log("[outline] polishing …")
    cleaned = polish_outline(cleaned, src_for_polish)
    cleaned = allocate_details_across_slides(cleaned, src_for_polish)

    log(f"[outline] done. slides={len(cleaned)}")
    return {"deck_title": deck_title or "Presentation", "slides": cleaned}


# --- Build PPT ---
class BuildIn(BaseModel):
    outline_json: Dict[str, Any] = Field(..., description="Slide outline JSON")
    path: str = Field(default="output.pptx", description="Output path for .pptx")

def build_ppt(outline_json: Dict[str, Any], path: str = "output.pptx") -> str:
    prs = Presentation()

    def clean_line(s: str) -> str:
        return re.sub(r"^[\s\*\-•·]+", "", s or "").strip()

    slides = outline_json.get("slides", []) or []
    deck_title = clean_line(outline_json.get("deck_title", ""))

    if not slides:
        slides = [{"title": deck_title or "Welcome", "subtitle": "", "details": []},
                  {"title": "Conclusion", "subtitle": "", "details": ["(add point)", "(add point)", "(add point)"]},
                  {"title": "Thanks for Watching", "subtitle": "", "details": []}]

    # Prefer deck_title for filename; fallback to slide 1/2
    name_source = deck_title or clean_line(slides[0].get("title", "")) or (clean_line(slides[1].get("title", "")) if len(slides) > 1 else "Presentation")
    safe_title = re.sub(r"[^\w\-]+", "_", name_source).strip("_") or "Presentation"
    if not path or path == "output.pptx":
        path = f"{safe_title}.pptx"
        base, ext = os.path.splitext(path)
        i = 2
        while os.path.exists(path):
            path = f"{base}_{i}{ext}"
            i += 1

    total = len(slides)

    for idx, sl in enumerate(slides, start=1):
        title = clean_line(sl.get("title", f"Slide {idx}"))
        subtitle = clean_line(sl.get("subtitle", ""))
        details = [clean_line(x) for x in (sl.get("details") or []) if clean_line(x)]

        is_welcome = (idx == 1)
        is_thanks  = (idx == total)
        is_concl   = (idx == max(2, total-1))

        if is_welcome:
            layout = prs.slide_layouts[0]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = title
            if len(s.placeholders) > 1:
                s.placeholders[1].text = subtitle
            else:
                box = s.shapes.add_textbox(
                    left=prs.slide_width//10,
                    top=prs.slide_height//2,
                    width=prs.slide_width*8//10,
                    height=prs.slide_height//5
                )
                box.text_frame.text = subtitle or ""
            continue

        layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = title
        body = s.placeholders[1]

        if is_thanks:
            # Thanks slide should have no bullets
            body.text = subtitle or ""
            continue

        # Non-thanks slides: enforce at least 3 bullets on regular and conclusion
        if not details:
            details = ["(add point)", "(add point)", "(add point)"]
        elif (not is_welcome) and (not is_thanks) and len(details) < 3:
            details += ["(add point)"] * (3 - len(details))

        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.text = details[0]
        for d in details[1:]:
            p = tf.add_paragraph()
            p.text = d

    prs.save(path)
    return os.path.abspath(path)


Build = StructuredTool.from_function(
    func=build_ppt,
    name="build_ppt",
    description=(
        "Render a PowerPoint (.pptx) from slide_outline_json. Final step only."
    ),
    args_schema=BuildIn,
)

Outline = StructuredTool.from_function(
    func=slide_outline_json,
    name="slide_outline_json",
    description="Generate slide outline JSON (title, subtitle, details) from summarized notes",
    args_schema=OutlineIn,
)

TOOLS: Dict[str, StructuredTool] = {
    "read_pdf": ReadPDF,
    "summarize_text": Summarize,
    "slide_outline_json": Outline,
    "build_ppt": Build,
}

# =============================================================
# 2) LangGraph State & Nodes
# =============================================================

class AgentState(TypedDict, total=False):
    messages: List[Dict[str, str]]         # chat history [{role, content}]
    pdf_text: Optional[str]
    notes: Optional[str]
    outline: Optional[Dict[str, Any]]
    ppt_path: Optional[str]
    next_action: Optional[str]
    next_args: Optional[Dict[str, Any]]
    auto: bool                             # if True, don't ask before building

# Helper to get last user message
def last_user_text(state: AgentState) -> str:
    for m in reversed(state.get("messages", [])):
        if m["role"] == "user":
            return m["content"]
    return ""

BAD_PHRASE_PAT = re.compile(
    r"(?i)\b(key takeaways include:|conclusion/key takeaways|conclusion:|^[-*•]+\s*$|\*\*|^\*|\s\*$)"
)
SECTION_PREFIX_PAT = re.compile(
    r"(?i)^\s*(background/?definition|history/?origins|key concepts/?taxonomy|"
    r"real examples/?case studies|applications/?impact|health/?risks/?limitations|"
    r"best practices/?tips|recent trends/?stats)\s*[:\-]*\s*"
)

def fix_truncated_line(line: str, notes: str) -> str:
    s = line.rstrip()
    # heuristic: ends with a 1–3 letter fragment
    if re.search(r"[A-Za-z]{1,3}$", s) and not s.endswith((".", "!", "?")):
        # find a note sentence that begins with this prefix (case-insensitive)
        prefix = re.sub(r"\s+", " ", s[-40:]).strip().lower()
        for sent in split_sentences(notes):
            low = sent.lower()
            if low.startswith(prefix):
                return sent[:120].rstrip()
        # otherwise trim the dangling fragment
        s = re.sub(r"\s+[A-Za-z]{1,3}$", "", s).rstrip()
    return s

def clean_detail_line(s: str) -> str:
    s = re.sub(r"^[\s\*\-•·]+", "", (s or "").strip())
    s = SECTION_PREFIX_PAT.sub("", s)
    s = BAD_PHRASE_PAT.sub("", s)
    s = re.sub(r"\*\*+", "", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s

def jaccard(a: str, b: str) -> float:
    wa, wb = set(re.findall(r"[a-z0-9]+", a.lower())), set(re.findall(r"[a-z0-9]+", b.lower()))
    if not wa or not wb: return 0.0
    return len(wa & wb) / len(wa | wb)

def compress_details(lines: list[str], min_keep: int = 3, max_keep: int = 5, sim_thresh: float = 0.6) -> list[str]:
    out = []
    for ln in (clean_detail_line(x) for x in lines):
        if not ln: continue
        if any(jaccard(ln, y) >= sim_thresh for y in out):  # drop near-duplicates
            continue
        out.append(ln)
        if len(out) >= max_keep: break
    while len(out) < min_keep:
        out.append("(add point)")
    return out[:max_keep]

def polish_outline(cleaned: list[Dict[str, Any]], notes: str) -> list[Dict[str, Any]]:
    if not cleaned: return cleaned
    n = len(cleaned)
    # Fix middle slides (2..n-2)
    for i in range(1, max(1, n-2)):
        sl = cleaned[i]
        sl["title"] = repair_title(sl.get("title", ""), notes)
        details = sl.get("details") or []
        details = compress_details(details, min_keep=3, max_keep=5)
        if all("(add point)" in d for d in details):
            # synthesize from notes if still empty-ish
            details = synthesize_details_from_notes(sl["title"], notes, k=4)
            details = compress_details(details, min_keep=3, max_keep=5)
        sl["details"] = details

    # Ensure Conclusion (n-1)
    if n >= 2:
        concl = cleaned[-2]
        concl["title"] = "Conclusion"
        concl["details"] = compress_details(
            concl.get("details") or synthesize_details_from_notes("Key takeaways", notes, k=4),
            min_keep=3, max_keep=5
        )

    # Ensure Thanks (n)
    cleaned[-1]["title"] = "Thanks for Watching"
    cleaned[-1]["details"] = []

    # Welcome (1): keep subtitle tidy
    cleaned[0]["subtitle"] = re.sub(r"\s+", " ", (cleaned[0].get("subtitle") or f"An overview of {cleaned[0].get('title','')}")).strip()
    cleaned[0]["details"] = []
    return cleaned

def repair_title(t: str, notes: str) -> str:
    t = (t or "").strip()
    if not t or re.fullmatch(r"(have been|dogs highly|highly social|other dogs|dogs also)", t, re.I) or len(t.split()) < 2:
        # propose a better title from notes’ keyphrases
        cands = keyphrase_titles_from_notes(notes, 1)
        return cands[0] if cands else "Key Topic"
    # title-case and strip junk
    t = re.sub(r"\s+", " ", t)
    return t[0].upper() + t[1:]

def extract_pdf_path(text: str) -> str | None:
    """Return a plausible PDF path from a user message (Windows-friendly)."""
    # Match quoted paths: "D:\dir\file.pdf"
    m = re.search(r'"([^"]+\.pdf)"', text, re.I)
    if m:
        return m.group(1)
    # Match unquoted Windows path or bare name ending with .pdf
    m = re.search(r'([A-Za-z]:\\[^\s]+\.pdf)', text, re.I)
    if m:
        return m.group(1)
    m = re.search(r'(\S+\.pdf)\b', text, re.I)
    if m:
        return m.group(1)
    return None

def wants_summary_only_now(state: AgentState) -> bool:
    return "summary only" in (last_user_text(state).lower())

# Helper to append a message to state
def add_message(state: AgentState, role: str, content: str) -> AgentState:
    msgs = list(state.get("messages", []))
    msgs.append({"role": role, "content": content})
    state["messages"] = msgs
    return state

def _chunk_by_sentences(s: str, max_chars: int = 3000) -> list[str]:
    # Chunk into ~2–3k chars so each map call finishes fast
    sents = split_sentences(s)
    chunks, cur, total = [], [], 0
    for sent in sents:
        L = len(sent) + 1
        if total + L > max_chars and cur:
            chunks.append(" ".join(cur))
            cur, total = [], 0
        cur.append(sent)
        total += L
    if cur:
        chunks.append(" ".join(cur))
    return chunks

def _hash(s: str) -> str:
    import hashlib
    return hashlib.sha1(s.encode("utf-8", errors="ignore")).hexdigest()[:10]

def infer_topic_from_user_text(state: AgentState) -> str:
    t = last_user_text(state)
    t = re.sub(r"\b(make|create|build|slides?|ppt|presentation|about|on)\b", "", t, flags=re.I)
    t = re.sub(r"\s+", " ", t).strip()
    return t[:80].title() or "Presentation"

def split_sentences(text: str) -> list[str]:
    # very light splitter; avoids bringing in external deps
    parts = re.split(r'(?<=[.!?])\s+', text.strip())
    return [p.strip() for p in parts if p.strip()]

def normalize_line_global(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "").strip().lower())

def compress_details_unique(lines: list[str], global_seen: set[str],
                            min_keep: int = 3, max_keep: int = 5, sim_thresh: float = 0.6) -> list[str]:
    out = []
    for ln in (clean_detail_line(x) for x in lines):
        if not ln:
            continue
        norm = normalize_line_global(ln)
        if norm in global_seen:
            continue
        if any(jaccard(ln, y) >= sim_thresh for y in out):
            continue
        out.append(ln)
        global_seen.add(norm)
        if len(out) >= max_keep:
            break
    while len(out) < min_keep:
        out.append("(add point)")
    return out[:max_keep]

def allocate_details_across_slides(cleaned: list[Dict[str, Any]], notes: str) -> list[Dict[str, Any]]:
    if not cleaned: return cleaned
    n = len(cleaned)
    global_seen: set[str] = set()

    # Reserve welcome & thanks; seed global_seen with their lines (usually none)
    for idx in [0, n-1]:
        for ln in cleaned[idx].get("details", []) or []:
            global_seen.add(normalize_line_global(ln))

    # Walk through content slides in order, compressing uniquely and filling from notes with exclude=global_seen
    for i in range(1, max(1, n-1)):
        if i == n-1:  # thanks
            break
        sl = cleaned[i]
        if sl.get("title", "").lower() == "conclusion":
            # handle conclusion last in this loop
            continue

        # 1) Unique-compress what's there
        uniq = compress_details_unique(sl.get("details") or [], global_seen, min_keep=0, max_keep=5)

        # 2) If <3, synthesize from notes excluding global_seen
        if len(uniq) < 3:
            extra = synthesize_details_from_notes(sl.get("title", ""), notes, k=5-len(uniq), exclude=global_seen)
            for e in extra:
                norm = normalize_line_global(e)
                if norm in global_seen:
                    continue
                uniq.append(e)
                global_seen.add(norm)
                if len(uniq) >= 5:
                    break

        # 3) Guarantee 3–5
        while len(uniq) < 3:
            uniq.append("(add point)")
        sl["details"] = uniq[:5]

    # Finally ensure conclusion has unique lines too
    if n >= 2:
        concl = cleaned[-2]
        uniq = compress_details_unique(concl.get("details") or [], global_seen, min_keep=0, max_keep=5)
        if len(uniq) < 3:
            extra = synthesize_details_from_notes("Key takeaways", notes, k=5-len(uniq), exclude=global_seen)
            for e in extra:
                norm = normalize_line_global(e)
                if norm in global_seen:
                    continue
                uniq.append(e)
                global_seen.add(norm)
                if len(uniq) >= 5:
                    break
        while len(uniq) < 3:
            uniq.append("(add point)")

        # NEW: if placeholders remain, try LLM densify, then fallback synthesize without exclude
        if all("(add point)" in d for d in uniq):
            try:
                bt, bd = densify_slide_with_llm(sl.get("title",""), cleaned[0].get("title","Presentation"), notes, max_points=5)
                cand = [clean_detail_line(fix_truncated_line(x, notes)) for x in bd]
                uniq = [x for x in cand if x][:5] or uniq
            except Exception:
                pass
            if all("(add point)" in d for d in uniq):
                extra = synthesize_details_from_notes(sl.get("title",""), notes, k=5, exclude=set())  # <— no exclude now
                cand = [clean_detail_line(fix_truncated_line(x, notes)) for x in extra]
                uniq = [x for x in cand if x][:5] or uniq

        sl["details"] = uniq[:5]
    return cleaned


def keyphrase_titles_from_notes(notes: str, k: int) -> list[str]:
    text = re.sub(r"[^a-z0-9\s]", " ", notes.lower())
    words = [w for w in text.split() if 2 <= len(w) <= 20]
    if not words:
        return []
    # stopwords-lite
    stop = set("a an the of to for and or in on with at by from as is are be was were this that these those it its their our your you we us".split())
    words = [w for w in words if w not in stop]
    # generate bi/tri-grams
    grams = []
    for n in (3, 2):
        for i in range(len(words)-n+1):
            grams.append(" ".join(words[i:i+n]))
    # count + rank
    freq: Dict[str,int] = {}
    for g in grams:
        freq[g] = freq.get(g, 0) + 1
    ranked = sorted(freq.items(), key=lambda x: (-x[1], -len(x[0])))
    titles = []
    for g, _ in ranked:
        # title case, de-dup by substring overlap
        cand = " ".join(w.capitalize() for w in g.split())
        if any(cand in t or t in cand for t in titles):
            continue
        titles.append(cand)
        if len(titles) >= k:
            break
    return titles

def is_generic_title(t: str) -> bool:
    return bool(GENERIC_TITLE_PAT.match((t or "").strip()))

def make_deck_title_from_notes(notes: str) -> str:
    # build a readable deck name from top keyphrases or first good sentence
    keys = keyphrase_titles_from_notes(notes, 1)
    if keys:
        return keys[0]
    sents = split_sentences(notes)
    if sents:
        return sents[0][:60]
    return "Presentation"

def tokenize(s: str) -> list[str]:
    return re.findall(r"[a-zA-Z0-9]+", (s or "").lower())

def synthesize_details_from_notes(title: str, notes: str, k: int = 4, exclude: Optional[set[str]] = None) -> list[str]:
    exclude = exclude or set()
    title_terms = set(w for w in tokenize(title) if len(w) > 2)
    sents = split_sentences(notes)
    if not sents:
        return ["(add point)"] * max(3, k)

    scored = []
    for s in sents:
        raw = s.strip()
        key = re.sub(r"\s+", " ", raw.lower())
        if key in exclude:
            continue
        ws = set(w for w in tokenize(raw) if len(w) > 2)
        overlap = len(ws & title_terms)
        length = len(raw)
        length_score = -abs(length - 90) / 90.0
        score = overlap * 3.0 + length_score
        scored.append((score, raw))

    scored.sort(reverse=True, key=lambda x: x[0])

    picked = []
    seen_local = set()
    for _, raw in scored:
        line = re.sub(r"^[\s\*\-•·]+", "", raw)
        norm = re.sub(r"\s+", " ", line.lower()).strip()
        if not norm or norm in seen_local or norm in exclude:
            continue
        # shorten nicely
        if len(line) > 100:
            cut = re.search(r"[;,:–—\-]\s", line[60:100])
            line = (line[:60 + cut.start()].strip() if cut else line[:90].rstrip())
        picked.append(line)
        seen_local.add(norm)
        if len(picked) >= k:
            break

    while len(picked) < max(3, k):
        picked.append("(add point)")
    return picked[: max(3, k)]

def extract_quant_facts(text: str, max_items: int = 20) -> list[str]:
    """
    Pull out quantitative tidbits (percentages, counts, sizes, model names) to preserve in summary/outline.
    """
    import re
    t = " " + re.sub(r"\s+", " ", text) + " "
    patterns = [
        r"\b\d{1,3}(?:\.\d+)?\s?%",                                # 22%, 7.5%
        r"\b(?:\d{1,3}(?:,\d{3})+|\d+)\b\s*(?:samples?|tokens?|params?|parameters|examples?|instances?|records?|docs?)\b",
        r"\b(?:\d{1,3}(?:\.\d+)?\s?(?:M|B|G|K)\+?)\s*(?:params?|parameters|tokens?)\b",  # 7B params, 50M tokens
        r"\b(?:top|avg|mean|median)\s*[- ]?\s*\d{1,3}(?:\.\d+)?\b", # top-1 73.2
        r"\b\d+\s+(?:models?|datasets?|benchmarks?|tasks?|epochs|layers|heads)\b",      # 7 models, 12 epochs
        r"\b[A-Z][A-Za-z0-9\-\+]{2,}\s?(?:[0-9]{0,3}[Bb])?\b",     # model names like GPT-4, Llama-3, Mistral-7B
        r"\bAUROC\s*\d{1,3}(?:\.\d+)?\b",                          # AUROC 88.4
        r"\b(?:improv(?:e|ement)|gain|boost)\s*(?:of|by)?\s*\d{1,3}(?:\.\d+)?\s?%\b",
    ]
    hits = []
    for pat in patterns:
        for m in re.finditer(pat, t, flags=re.I):
            frag = t[max(0, m.start()-40): m.end()+40].strip()
            frag = re.sub(r"\s+", " ", frag)
            hits.append(frag)
    # de-dupe by lowercased text
    out, seen = [], set()
    for h in hits:
        k = h.lower()
        if k in seen: 
            continue
        seen.add(k)
        out.append(h)
        if len(out) >= max_items:
            break
    return out

# --- Planner Node: LLM decides next action (which tool/chat/finish) ---
def planner_node(state: AgentState) -> AgentState:
    llm = get_llm_singleton(temp=0.0, num_predict=256)
    user_text = last_user_text(state).strip()
    text_l = user_text.lower()
    state["next_action"] = None
    state["next_args"] = {}

    if state.get("ppt_path"):
        state["next_action"] = "finish"; state["next_args"] = {}; return state

    def extract_slide_count(s: str, default: int = 11) -> int:
        m = re.search(r"(\d+)\s*(slides?|ppt)", s.lower())
        return int(m.group(1)) if m else default

    wants_ppt = any(k in text_l for k in [
        "ppt","slides","slide deck","powerpoint","power point","presentation","presentaion"
    ])
    wants_summary_only = ("summary only" in text_l) or (("summary" in text_l) and not wants_ppt)
    mentions_pdf = ".pdf" in text_l
    asks_regen_outline = "regenerate outline" in text_l or "re-generate outline" in text_l or "redo outline" in text_l
    asks_build_now = "build slides now" in text_l or text_l == "build slides" or "build now" in text_l

    slide_target = extract_slide_count(text_l, default=11)
    state["last_slide_target"] = slide_target
    pdf_path = extract_pdf_path(user_text) if mentions_pdf else None

    # --- SUMMARY-ONLY path: read → summarize → finish (no loops)
    if wants_summary_only and mentions_pdf and pdf_path and not state.get("pdf_text"):
        state["next_action"] = "read_pdf"; state["next_args"] = {"path": pdf_path}; return state
    if wants_summary_only and state.get("pdf_text") and not state.get("notes"):
        state["next_action"] = "summarize_text"; state["next_args"] = {"text": state["pdf_text"], "style": "exec"}; return state
    if wants_summary_only and state.get("notes"):
        state["next_action"] = "finish"; state["next_args"] = {}; return state

    # --- Controls
    if asks_regen_outline and state.get("notes"):
        state["next_action"] = "slide_outline_json"; state["next_args"] = {"notes": state["notes"], "slide_target": slide_target}; return state

    # --- BUILD SLIDES pressed (this enforces summarize → outline → build)
    if asks_build_now:
        # a) Already have outline → just build
        if state.get("outline"):
            state["next_action"] = "build_ppt"
            state["next_args"] = {"outline_json": state["outline"], "path": "output.pptx"}
            return state
        # b) Have notes (summary) → outline next, then auto-build
        if state.get("notes"):
            state["auto"] = True
            state["next_action"] = "slide_outline_json"
            state["next_args"] = {"notes": state["notes"], "slide_target": slide_target}
            return state
        # c) Have pdf_text but no notes → summarize first, then auto-chain
        if state.get("pdf_text"):
            state["auto"] = True
            state["next_action"] = "summarize_text"
            state["next_args"] = {"text": state["pdf_text"], "style": "slides"}
            return state
        # d) Nothing read → topic outline then auto-build
        topic = infer_topic_from_user_text(state)
        state["auto"] = True
        state["next_action"] = "slide_outline_json"
        state["next_args"] = {"topic": topic, "slide_target": slide_target}
        return state

    # --- Topic-only PPT (no PDF): go straight to outline
    if wants_ppt and not mentions_pdf:
        topic = infer_topic_from_user_text(state)
        state["next_action"] = "slide_outline_json"; state["next_args"] = {"topic": topic, "slide_target": slide_target}; return state

    # --- PDF + PPT (planned over turns)
    if wants_ppt and mentions_pdf and pdf_path and not state.get("pdf_text"):
        state["next_action"] = "read_pdf"; state["next_args"] = {"path": pdf_path}; return state
    if wants_ppt and state.get("pdf_text") and not state.get("notes"):
        state["next_action"] = "summarize_text"; state["next_args"] = {"text": state["pdf_text"], "style": "slides"}; return state
    if wants_ppt and state.get("notes") and not state.get("outline"):
        state["next_action"] = "slide_outline_json"; state["next_args"] = {"notes": state["notes"], "slide_target": slide_target}; return state
    if wants_ppt and state.get("outline"):
        state["next_action"] = "build_ppt"; state["next_args"] = {"outline_json": state["outline"], "path": "output.pptx"}; return state

    # default chat
    state["next_action"] = "chat"; state["next_args"] = {}; return state

# --- Tool Executor Node: run the chosen tool (if any) ---
def tool_node(state: AgentState) -> AgentState:
    import os, re, json, datetime
    action = (state.get("next_action") or "").strip()
    args: dict = (state.get("next_args") or {}) if isinstance(state.get("next_args"), dict) else {}

    if not action:
        state = add_message(state, "assistant", "No action planned. Tell me what to do (e.g., `summary only for <pdf>`).")
        state["next_action"] = "finish"; state["next_args"] = {}; return state

    # Prefill
    if action == "summarize_text":
        if not args.get("text"):
            seed = state.get("pdf_text") or f"{last_user_text(state)}\n\nCreate comprehensive notes about this topic for slides."
            args["text"] = seed  # no clamp
        args.setdefault("style", "slides")
        args.setdefault("target_words", 350)

    elif action == "slide_outline_json":
        if not args.get("notes") and not args.get("topic"):
            args["topic"] = infer_topic_from_user_text(state)
        args.setdefault("slide_target", 11)
        args.setdefault("verbose", True)
        
    elif action == "build_ppt":
        if not args.get("outline_json"):
            if state.get("outline"):
                args["outline_json"] = state["outline"]
            else:
                state = add_message(state, "assistant", "I need an outline before building slides. Say `regenerate outline` or ask me to create an outline.")
                state["next_action"] = "finish"; state["next_args"] = {}; return state
        if not args.get("path"):
            deck_title = (args["outline_json"] or {}).get("deck_title") or "Presentation"
            safe = re.sub(r"[^\w\-]+", "_", deck_title).strip("_") or "Presentation"
            ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            args["path"] = os.path.abspath(f"{safe}_{ts}.pptx")

    # Run tool
    if action in TOOLS:
        tool = TOOLS[action]
        try:
            result = tool.invoke(args)
        except Exception as e:
            msg = str(e)
            state = add_message(state, "assistant", f"❌ Tool `{action}` failed: {e}")
            state["next_action"] = "finish"; state["next_args"] = {}; return state
    else:
        result = None

    # Post-process + auto-chain
    if action == "read_pdf":
        text = result or ""
        state["pdf_text"] = text
        state = add_message(state, "assistant", f"Read PDF ✔ (text cached, {len(text)} chars). You can now say `summary only` or `build slides`.")

    elif action == "summarize_text":
        summary = (result or "").strip()
        state["notes"] = summary

        if state.get("auto"): 
            state["auto"] = False
            try:
                target = int(state.get("last_slide_target", 11))
                outline = TOOLS["slide_outline_json"].invoke({"notes": summary, "slide_target": target})
            except Exception as e:
                state = add_message(state, "assistant", f"Summary ready ✔\n\n{summary}\n\n(Outline failed during auto-chain: {e})")
                state["next_action"] = "finish"; state["next_args"] = {}; return state

            state["outline"] = outline
            deck_title = (outline or {}).get("deck_title") or "Presentation"
            safe = re.sub(r"[^\w\-]+", "_", deck_title).strip("_") or "Presentation"
            ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            out_path = os.path.abspath(f"{safe}_{ts}.pptx")

            try:
                out_file = TOOLS["build_ppt"].invoke({"outline_json": outline, "path": out_path})
                state["ppt_path"] = out_file
                state = add_message(state, "assistant", f"Summary ready ✔\n\n{summary}\n\nPPT built ✔ -> {out_file}")
            except Exception as e:
                pretty = json.dumps(outline, indent=2, ensure_ascii=False)
                state = add_message(state, "assistant", "Summary ready ✔\n\n" + summary +
                                    "\n\nSlide outline JSON prepared — review below:\n\n```json\n" + pretty + "\n```\n"
                                    f"(Auto-build failed: {e})\nSay **'build slides'** to render.")
        else:
            state = add_message(state, "assistant", "Summary ready ✔\n\n" + summary)

    elif action == "slide_outline_json":
        outline = result or {}
        state["outline"] = outline
        pretty = json.dumps(outline, indent=2, ensure_ascii=False)
        try:
            with open("outline_debug.json", "w", encoding="utf-8") as f:
                f.write(pretty)
            saved_note = " (saved to outline_debug.json)"
        except Exception:
            saved_note = ""

        if state.get("auto"):  # user said "build slides" → build now
            state["auto"] = False
            deck_title = (outline or {}).get("deck_title") or "Presentation"
            safe = re.sub(r"[^\w\-]+", "_", deck_title).strip("_") or "Presentation"
            ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            out_path = os.path.abspath(f"{safe}_{ts}.pptx")
            try:
                out_file = TOOLS["build_ppt"].invoke({"outline_json": outline, "path": out_path})
                state["ppt_path"] = out_file
                state = add_message(state, "assistant", f"Slide outline JSON prepared ✔{saved_note}\n\nPPT built ✔ -> {out_file}")
            except Exception as e:
                state = add_message(state, "assistant", "Slide outline JSON prepared ✔" + saved_note +
                                    " — review below:\n\n```json\n" + pretty + "\n```\n"
                                    f"(Auto-build failed: {e})\nSay **'build slides'** to render.")
        else:
            state = add_message(state, "assistant",
                "Slide outline JSON prepared ✔" + saved_note + " — review below:\n\n```json\n" + pretty + "\n```\n"
                "Say **'build slides'** to render this outline, or **'regenerate outline'** to try again."
            )

    elif action == "build_ppt":
        out_path = result if (isinstance(result, str) and result.lower().endswith(".pptx")) else args.get("path")
        state["ppt_path"] = out_path
        state = add_message(state, "assistant", f"PPT built ✔ -> {out_path}")

    elif action == "chat":
        state = add_message(state, "assistant", "How can I help next? For example:\n - `summary only for \"D:/path/paper.pdf\"`\n - `read \"D:/path/paper.pdf\" pages 1-5 and summarize`\n - `make 10 slides ppt about diffusion models`")

    state["next_action"] = "finish"; state["next_args"] = {}
    return state

# --- Chat Node: simple chat response (no tools) ---
def chat_node(state: AgentState) -> AgentState:
    llm = get_llm_singleton(temp=0.0, num_predict=320)

    # Short system guidance to keep answers crisp
    sys = SystemMessage(content=(
        "You are a helpful research→slides assistant. "
        "Answer briefly and clearly. If the user asks for PPTs or summaries, "
        "you may propose a plan, but do not call tools from this node."
    ))

    # Include a small window of history (last few turns)
    history = [sys]
    for m in state.get("messages", [])[-6:]:
        if m["role"] == "user":
            history.append(HumanMessage(content=m["content"]))
        else:
            history.append(AIMessage(content=m["content"]))

    reply = llm.invoke(history).content
    state = add_message(state, "assistant", reply)

    # Clear any stale next_* so the router doesn't reuse them
    state["next_action"] = None
    state["next_args"] = {}
    return state

# --- Decider for graph edges ---

def choose_route(state: AgentState) -> Literal["tool", "chat", "finish"]:
    action = (state.get("next_action") or "chat").lower()
    if action == "finish":
        return "finish"
    if action in TOOLS:
        return "tool"
    if action == "chat":
        return "chat"
    return "chat"


# =============================================================
# 3) Build Graph Application
# =============================================================

def build_app():
    g = StateGraph(AgentState)
    g.add_node("plan", planner_node)
    g.add_node("tool", tool_node)
    g.add_node("chat", chat_node)

    g.set_entry_point("plan")

    g.add_conditional_edges("plan", choose_route, {
        "tool": "tool",
        "chat": "chat",
        "finish": END,
    })
    # Important: end the turn after any tool to avoid auto-loops
    g.add_edge("tool", END)
    g.add_edge("chat", END)
    return g.compile()

# =============================================================
# 4) Simple CLI loop
# =============================================================

if __name__ == "__main__":
    hf_id = os.environ.get("HF_MODEL_ID", "Qwen/Qwen2.5-7B-Instruct")
    print(f"🧠 Using HF model: {hf_id}")

    app = build_app()

    state: AgentState = {"messages": []}

    print("\n📽️ LangGraph PPT Agent ready. Type your request (or 'exit').")
    print("Examples:\n - summary only for D:/papers/attention.pdf\n - make slides about diffusion models (10 slides)\n - read D:/papers/paper.pdf pages 1-3 and summarize\n - build slides now\n")

    steps = 0
    while True:
        try:
            user = input("You> ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\nBye!")
            break
        if user.lower() in {"exit", "quit"}:
            print("Bye!")
            break
        state = add_message(state, "user", user)
        state = app.invoke(state)   # exactly once
        if state.get("messages"):
            last = state["messages"][-1]
            if last["role"] == "assistant":
                print("Assistant>", last["content"])
