#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ppt_agent_langgraph_ollama_v8.py — LangGraph ReAct agent using Ollama
- Fix: add explicit docstrings to all @tool functions to satisfy LangChain's StructuredTool requirements.
- Keeps v7 improvements: stronger guidance, fallback build, verbose tracing, recursion guard, debounce.
"""
from __future__ import annotations

import argparse
import json
import os
import re
import sys
import time
import hashlib
from dataclasses import dataclass, asdict
from typing import Any, Dict, List, Optional

from langchain_ollama import ChatOllama
from langchain_core.messages import HumanMessage, SystemMessage, AIMessage, ToolMessage
from langchain_core.tools import tool
from langgraph.prebuilt import create_react_agent

# Optional deps used by tools
def _safe_imports():
    errors = []
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception as e:
        errors.append(("pypdf", e))
        PdfReader = None  # type: ignore
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        errors.append(("python-pptx", e))
        Presentation = None  # type: ignore
    return {"PdfReader": PdfReader, "Presentation": Presentation, "errors": errors}

DEPS = _safe_imports()

DEFAULTS = {
    "MODEL": os.getenv("OLLAMA_MODEL", "gpt-oss:20b"),
    "STATE_PATH": os.path.abspath("./agent_state.json"),
    "OUTPUT_DIR": os.path.abspath("./decks"),
    "DEFAULT_SLIDES": int(os.getenv("DEFAULT_SLIDES", "11")),
}

SYSTEM_PROMPT = """You are a presentation agent that MUST use tools to accomplish tasks.
Policy:
- For any request that matches ‘build .*ppt.* about <TOPIC>’ (allow extra words like Hello), you MUST:
  1) Call slide_outline_json(topic=<TOPIC>, slide_target=DEFAULT_SLIDES) exactly once;
  2) Then call build_ppt(outline_json=that JSON) exactly once;
  3) Then give a short final answer with the file path.
- If the user provides a file path, use: read_pdf -> summarize_text -> slide_outline_json -> build_ppt.
- For “tell me about the ppt” / “what’s in the slides”, call describe_latest_deck and DO NOT rebuild.
- Use each tool at most once per user request unless NEW information is provided.
- Maintain continuity via the tools (they persist topic, last_outline, ppt_path).
- Keep answers short and clear. Stop tool use when done.

Notes:
- DEFAULT_SLIDES = {ds} (you may override if user specifies a different number).
""".format(ds=DEFAULTS["DEFAULT_SLIDES"])

# ---------------- Persistence ----------------
@dataclass
class AgentStatePersisted:
    topic: Optional[str] = None
    last_outline: Optional[str] = None  # JSON string
    ppt_path: Optional[str] = None

    @classmethod
    def load(cls, path: str) -> "AgentStatePersisted":
        if not os.path.exists(path):
            return cls()
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return cls(**data)

    def save(self, path: str) -> None:
        with open(path, "w", encoding="utf-8") as f:
            json.dump(asdict(self), f, ensure_ascii=False, indent=2)


def _load_state() -> AgentStatePersisted:
    return AgentStatePersisted.load(DEFAULTS["STATE_PATH"])


def _save_state(s: AgentStatePersisted) -> None:
    s.save(DEFAULTS["STATE_PATH"])


# ---------------- Debounce helpers ----------------
_RECENT: Dict[str, float] = {}
def _sig(obj: Any) -> str:
    return hashlib.sha1(json.dumps(obj, sort_keys=True).encode("utf-8")).hexdigest()

def _recent_once(key: str, ttl: int = 90) -> bool:
    """Return True if this key has been seen recently (within ttl)."""
    now = time.time()
    last = _RECENT.get(key)
    _RECENT[key] = now
    # cleanup
    for k,v in list(_RECENT.items()):
        if now - v > ttl:
            _RECENT.pop(k, None)
    return last is not None and now - last < ttl


# ---------------- Tools ----------------
@tool
def read_pdf(path: str) -> dict:
    """Read a PDF or plaintext (.txt/.md) file from disk and return the raw text.
    Args:
        path: Absolute or relative path to a .pdf, .txt, or .md file.
    Returns:
        dict with keys: status, path, text or error.
    """
    PdfReader = DEPS["PdfReader"]
    if not os.path.exists(path):
        return {"status": "error", "error": f"File not found: {path}"}
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext in (".txt", ".md"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            return {"status": "ok", "path": path, "text": text[:2_000_000]}
        if ext == ".pdf":
            if PdfReader is None:
                return {"status": "error", "error": "pypdf not installed. pip install pypdf"}
            reader = PdfReader(path)
            chunks = []
            for page in reader.pages:
                try:
                    chunks.append(page.extract_text() or "")
                except Exception:
                    pass
            text = "\n".join(chunks)
            return {"status": "ok", "path": path, "text": text[:2_000_000]}
        return {"status": "error", "error": f"Unsupported file type: {ext}"}
    except Exception as e:
        return {"status": "error", "error": f"{type(e).__name__}: {e}"}


@tool
def summarize_text(text: str, max_words: int = 400) -> dict:
    """Create a short extractive summary from long text (no LLM cost).
    Args:
        text: Raw text to summarize.
        max_words: Approximate maximum words to keep across extracted sentences.
    Returns:
        dict with keys: status, summary.
    """
    import textwrap
    words_left = max_words
    out: List[str] = []
    for para in re.split(r"\n{2,}", text):
        para = para.strip()
        if not para:
            continue
        sent = re.split(r"(?<=[.!?])\s+", para)[0]
        w = sent.split()
        if not w:
            continue
        if len(w) > words_left:
            out.append(" ".join(w[:words_left]) + "…")
            words_left = 0
            break
        out.append(sent)
        words_left -= len(w)
        if words_left <= 0:
            break
    summary = "\n".join(textwrap.wrap(" ".join(out), width=100))
    st = _load_state()
    if summary:
        st.last_outline = None  # invalidate outline since notes changed
        _save_state(st)
    return {"status": "ok", "summary": summary}


def _ollama_json_outline(topic: Optional[str], notes: Optional[str], slide_target: int, model: str) -> str:
    """Internal helper: call ChatOllama to produce a strict-JSON slide outline; returns JSON string."""
    llm = ChatOllama(model=model, temperature=0.2)
    sys = f"""You output STRICT JSON for slide outlines only.
Schema:
{{
  "deck_title": string,
  "slides": [{{"title": string, "bullets": [string, ...]}}, ...]
}}
Rules:
- Exactly {slide_target} slides (1 title + {slide_target-1} content).
- Bullets are ≤12 words and non-redundant.
- Output JSON only. No code fences, no prose outside JSON.
"""
    parts = ["Create a deck outline."]
    if topic:
        parts.append(f"Topic: {topic}")
    if notes:
        parts.append(f"Use these notes:\n{notes[:8000]}")
    prompt = "\n".join(parts)
    resp = llm.invoke([SystemMessage(content=sys), HumanMessage(content=prompt)])
    content = (resp.content or "").strip()
    content = re.sub(r"^```json\s*|\s*```$", "", content, flags=re.S)
    try:
        json.loads(content)
    except Exception:
        content = json.dumps({
            "deck_title": topic or "Presentation",
            "slides": [{"title": topic or "Overview", "bullets": ["Outline generation failed; please retry."]}]
        }, ensure_ascii=False, indent=2)
    return content


_OUTLINE_CACHE: Dict[str, str] = {}

@tool
def slide_outline_json(topic: Optional[str] = None, notes: Optional[str] = None, slide_target: int = DEFAULTS["DEFAULT_SLIDES"], model: Optional[str] = None) -> dict:
    """Create a JSON outline for a PPT (deck_title + slides[title, bullets]) using the LLM.
    Args:
        topic: Topic for the presentation.
        notes: Optional notes to incorporate.
        slide_target: Desired number of slides including the title slide.
        model: Optional override of the Ollama model name.
    Returns:
        dict with keys: status ('ok' or 'cached'), outline_json (str).
    """
    model = model or DEFAULTS["MODEL"]
    key = _sig({"k":"outline","t":topic or "", "n":hashlib.sha1((notes or "").encode()).hexdigest(), "s":slide_target, "m":model})
    if _recent_once(key, ttl=90) and key in _OUTLINE_CACHE:
        content = _OUTLINE_CACHE[key]
        st = _load_state()
        st.last_outline = content
        if topic:
            try:
                data = json.loads(content)
                st.topic = data.get("deck_title") or topic
            except Exception:
                st.topic = topic
        _save_state(st)
        return {"status": "cached", "outline_json": content}

    content = _ollama_json_outline(topic, notes, slide_target, model)
    _OUTLINE_CACHE[key] = content
    st = _load_state()
    st.last_outline = content
    if topic:
        try:
            data = json.loads(content)
            st.topic = data.get("deck_title") or topic
        except Exception:
            st.topic = topic
    _save_state(st)
    return {"status": "ok", "outline_json": content}


@tool
def build_ppt(outline_json: str, output_dir: Optional[str] = None) -> dict:
    """Render a PPTX file from a slide outline JSON string.
    Args:
        outline_json: Strict JSON string with 'deck_title' and 'slides' entries.
        output_dir: Optional directory to write the PPTX file into.
    Returns:
        dict with keys: status, path (absolute), or error.
    """
    Presentation = DEPS["Presentation"]
    if Presentation is None:
        return {"status": "error", "error": "python-pptx not installed. pip install python-pptx"}

    sig = _sig({"k":"build","outline": outline_json})
    now = time.time()
    if _recent_once(sig, ttl=120):
        st = _load_state()
        return {"status": "skipped", "reason": "recently_built", "path": st.ppt_path}

    try:
        data = json.loads(outline_json)
    except Exception:
        data = {"deck_title": "Presentation", "slides": [{"title": "Overview", "bullets": ["(Invalid JSON)"]}]}

    os.makedirs(output_dir or DEFAULTS["OUTPUT_DIR"], exist_ok=True)
    safe_title = re.sub(r"[^a-zA-Z0-9._-]+", "_", (data.get("deck_title") or "Presentation"))[:80] or "Presentation"
    fname = f"{safe_title}_{int(time.time())}.pptx"
    path = os.path.join(output_dir or DEFAULTS["OUTPUT_DIR"], fname)

    prs = Presentation()
    title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[1]
    bullet_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    slides = data.get("slides") or []
    for idx, s in enumerate(slides):
        layout = title_layout if idx == 0 else bullet_layout
        slide = prs.slides.add_slide(layout)
        shapes = slide.shapes
        title = shapes.title if hasattr(shapes, "title") else None
        if title:
            title.text = s.get("title") or f"Slide {idx+1}"
        body = None
        for shp in shapes:
            if shp.is_placeholder and getattr(shp, "placeholder_format", None) and shp.placeholder_format.type != 1:
                body = shp.text_frame
                break
        bullets = s.get("bullets") or []
        if body:
            body.clear()
            if bullets:
                body.text = bullets[0]
                for b in bullets[1:]:
                    p = body.add_paragraph(); p.text = b; p.level = 0

    from pathlib import Path as _P
    prs.save(path)
    st = _load_state()
    st.ppt_path = str(_P(path).resolve())
    _save_state(st)
    return {"status": "ok", "path": st.ppt_path}


@tool
def describe_latest_deck() -> dict:
    """Describe the most recently built deck without rebuilding it.
    Returns:
        dict with keys: status ('ok' or 'no_deck'), path, topic, headings (list of slide titles).
    """
    st = _load_state()
    if not st.ppt_path:
        return {"status": "no_deck"}
    headings: List[str] = []
    if st.last_outline:
        try:
            data = json.loads(st.last_outline)
            for s in (data.get("slides") or [])[:8]:
                headings.append(s.get("title") or s.get("heading") or "Slide")
        except Exception:
            pass
    return {"status": "ok", "path": st.ppt_path, "topic": st.topic, "headings": headings}


TOOLS = [read_pdf, summarize_text, slide_outline_json, build_ppt, describe_latest_deck]


def ensure_dirs():
    os.makedirs(DEFAULTS["OUTPUT_DIR"], exist_ok=True)


def _pretty(obj: Any, max_len: int = 240) -> str:
    try:
        s = json.dumps(obj, ensure_ascii=False)
    except Exception:
        s = str(obj)
    return (s[:max_len] + "…") if len(s) > max_len else s


def _extract_topic_from_text(text: str) -> Optional[str]:
    """Heuristic: capture topic after the word 'about'."""
    m = re.search(r"about\s+(.+)$", text, flags=re.I)
    if m:
        return m.group(1).strip(" .!?\"'")
    return None


def main():
    parser = argparse.ArgumentParser(description="LangGraph ReAct agent (Ollama)")
    parser.add_argument("--model", default=DEFAULTS["MODEL"], help="Ollama model, e.g. gpt-oss:20b")
    parser.add_argument("--once", action="store_true", help="Read a single prompt from stdin then exit")
    parser.add_argument("--recursion", type=int, default=12, help="LangGraph recursion_limit per turn (default 12)")
    parser.add_argument("--verbose", action="store_true", help="Print step-by-step agent trace")
    parser.add_argument("--slides", type=int, default=DEFAULTS["DEFAULT_SLIDES"], help="Default slide count if user doesn't specify")
    args = parser.parse_args()

    # Reflect default slides into the prompt dynamically
    global SYSTEM_PROMPT
    SYSTEM_PROMPT = SYSTEM_PROMPT.replace("DEFAULT_SLIDES = "+str(DEFAULTS["DEFAULT_SLIDES"]), f"DEFAULT_SLIDES = {args.slides}")

    ensure_dirs()

    llm = ChatOllama(model=args.model, temperature=0.2)
    graph = create_react_agent(llm, tools=TOOLS)

    # Maintain history across turns; inject SystemMessage once
    history: List[Any] = [SystemMessage(content=SYSTEM_PROMPT)]

    def run_turn(user_text: str) -> str:
        nonlocal history
        history += [HumanMessage(content=user_text)]
        final = None
        step = 1
        for event in graph.stream({"messages": history}, stream_mode="values", config={"recursion_limit": args.recursion}):
            final = event
            if not args.verbose:
                continue
            msgs = final.get("messages", [])
            if not msgs:
                continue
            last = msgs[-1]
            if isinstance(last, AIMessage):
                if getattr(last, "tool_calls", None):
                    for tc in last.tool_calls:
                        name = getattr(tc, "name", "tool")
                        arguments = getattr(tc, "args", {})
                        print(f"[step {step}] agent → tool: {name} args={_pretty(arguments)}")
                        step += 1
                else:
                    text = (last.content or "").strip()
                    if text:
                        print(f"[step {step}] agent message: {_pretty(text, 400)}")
                        step += 1
            elif isinstance(last, ToolMessage):
                name = getattr(last, "name", "tool")
                try:
                    content = json.loads(last.content) if isinstance(last.content, str) else last.content
                except Exception:
                    content = last.content
                brief = {}
                if isinstance(content, dict):
                    for k in ("status", "path", "topic", "headings", "reason", "outline_json"):
                        if k in content:
                            brief[k] = content[k] if k != "outline_json" else "(json outline)"
                else:
                    brief = {"content": content}
                print(f"[step {step}] tool → agent: {name} result={_pretty(brief)}")
                step += 1

        # Fallback: if model didn't issue tool calls but asked to build
        if final and "messages" in final and final["messages"]:
            history = final["messages"]
            if not history or not isinstance(history[0], SystemMessage):
                history = [SystemMessage(content=SYSTEM_PROMPT)] + history

            text_lower = user_text.lower()
            if ("ppt" in text_lower or "presentation" in text_lower or "slides" in text_lower) and "about" in text_lower:
                last_ai = next((m for m in reversed(history) if isinstance(m, AIMessage)), None)
                issued_tool_call = bool(getattr(last_ai, "tool_calls", None))
                if not issued_tool_call:
                    topic = _extract_topic_from_text(user_text) or "Presentation"
                    # Fallback: outline -> build
                    result_outline = slide_outline_json.invoke({"topic": topic, "slide_target": args.slides})
                    if args.verbose:
                        print(f"[fallback] agent → tool: slide_outline_json args={_pretty({'topic': topic, 'slide_target': args.slides})}")
                        print(f"[fallback] tool → agent: slide_outline_json result={_pretty({k:('(json outline)' if k=='outline_json' else v) for k,v in result_outline.items()})}")
                    result_build = build_ppt.invoke({"outline_json": result_outline.get("outline_json", "{}")})
                    if args.verbose:
                        print(f"[fallback] agent → tool: build_ppt args={{'outline_json': '(json outline)'}}")
                        print(f"[fallback] tool → agent: build_ppt result={_pretty(result_build)}")
                    path = result_build.get("path")
                    answer = f"I created your deck about “{topic}”. File: {path}" if path else "I attempted to build the deck but didn’t get a file path."
                    history += [AIMessage(content=answer)]
                    return answer

            last_ai = next((m for m in reversed(history) if isinstance(m, AIMessage)), None)
            return (last_ai.content if last_ai else "(no response)") or ""

        return "(no response)"

    if args.once:
        user_text = sys.stdin.read().strip()
        reply = run_turn(user_text)
        print(reply)
        return

    print("💡 ppt_agent (LangGraph ReAct + Ollama) — type your request. Ctrl+C to exit.")
    print(f"Model: {args.model}")
    while True:
        try:
            user_text = input("> ").strip()
            if not user_text:
                continue
            reply = run_turn(user_text)
            print(reply)
        except KeyboardInterrupt:
            print("\nBye.")
            break


if __name__ == "__main__":
    main()
