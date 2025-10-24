from __future__ import annotations

import argparse
import json
import os
import re
import sys
import time
import hashlib
from dataclasses import dataclass, asdict
from typing import Any, Dict, List, Optional,Union

# LangChain / LangGraph
from langchain_ollama import ChatOllama
from langchain_core.messages import HumanMessage, SystemMessage, AIMessage, ToolMessage
from langchain_core.tools import tool
from langgraph.prebuilt import create_react_agent

# ---------------- Optional deps used by tools ----------------
def _safe_imports():
    errors = []
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception as e:
        errors.append(("pypdf", e))
        PdfReader = None  # type: ignore
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        errors.append(("python-pptx", e))
        Presentation = None  # type: ignore
    return {"PdfReader": PdfReader, "Presentation": Presentation, "errors": errors}

DEPS = _safe_imports()

# ---------------- Defaults & persistence ----------------
DEFAULTS = {
    "MODEL": os.getenv("OLLAMA_MODEL", "gpt-oss:20b"),
    "STATE_PATH": os.path.abspath("./agent_state.json"),
    "OUTPUT_DIR": os.path.abspath("./decks"),
    "DEFAULT_SLIDES": int(os.getenv("DEFAULT_SLIDES", "11")),
}

SYSTEM_PROMPT = """You are a presentation agent that MUST use tools to accomplish tasks.
Policy:
- For any request that matches ‘build .*ppt.* about <TOPIC>’ (allow extra words like Hello), you MUST:
  1) Call slide_outline_json(topic=<TOPIC>, slide_target=DEFAULT_SLIDES) exactly once;
  2) Then call build_ppt(outline_json=that JSON) exactly once;
  3) Then give a short final answer with the file path.
- If the user provides a file path, use: read_pdf -> summarize_text -> slide_outline_json -> build_ppt.
- For “tell me about the ppt” / “what’s in the slides”, call describe_latest_deck and DO NOT rebuild.
- Use each tool at most once per user request unless NEW information is provided.
- Maintain continuity via the tools (they persist topic, last_outline, ppt_path).
- Keep answers short and clear. Stop tool use when done.

Notes:
- DEFAULT_SLIDES = {ds} (you may override if user specifies a different number).
""".format(ds=DEFAULTS["DEFAULT_SLIDES"])

@dataclass
class AgentStatePersisted:
    topic: Optional[str] = None
    last_outline: Optional[str] = None   # JSON string
    ppt_path: Optional[str] = None
    last_text: Optional[str] = None      # NEW: raw text from read_pdf
    last_summary: Optional[str] = None   # NEW: summary from summarize_text

    @classmethod
    def load(cls, path: str) -> "AgentStatePersisted":
        if not os.path.exists(path):
            return cls()
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return cls(**data)

    def save(self, path: str) -> None:
        with open(path, "w", encoding="utf-8") as f:
            json.dump(asdict(self), f, ensure_ascii=False, indent=2)

def _load_state() -> AgentStatePersisted:
    return AgentStatePersisted.load(DEFAULTS["STATE_PATH"])

def _save_state(s: AgentStatePersisted) -> None:
    s.save(DEFAULTS["STATE_PATH"])

# ---------------- Debounce helpers ----------------
_RECENT: Dict[str, float] = {}
def _sig(obj: Any) -> str:
    return hashlib.sha1(json.dumps(obj, sort_keys=True).encode("utf-8")).hexdigest()

def _recent_once(key: str, ttl: int = 90) -> bool:
    """Return True if this key has been seen recently (within ttl)."""
    now = time.time()
    last = _RECENT.get(key)
    _RECENT[key] = now
    # cleanup old
    for k, v in list(_RECENT.items()):
        if now - v > ttl:
            _RECENT.pop(k, None)
    return (last is not None) and (now - last < ttl)

# ---------------- Tools ----------------
@tool
def read_pdf(path: str) -> dict:
    """Read a PDF or plaintext (.txt/.md) file and return the raw text.
    Persists last_text and clears last_summary.
    """
    PdfReader = DEPS["PdfReader"]
    if not os.path.exists(path):
        return {"status": "error", "error": f"File not found: {path}"}
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext in (".txt", ".md"):
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif ext == ".pdf":
            if PdfReader is None:
                return {"status": "error", "error": "pypdf not installed. pip install pypdf"}
            reader = PdfReader(path)
            chunks = []
            for page in reader.pages:
                try:
                    chunks.append(page.extract_text() or "")
                except Exception:
                    pass
            text = "\n".join(chunks)
        else:
            return {"status": "error", "error": f"Unsupported file type: {ext}"}

        st = _load_state()
        st.last_text = text[:2_000_000]
        st.last_summary = None  # new text invalidates any previous summary
        _save_state(st)

        return {"status": "ok", "path": path, "text": st.last_text}
    except Exception as e:
        return {"status": "error", "error": f"{type(e).__name__}: {e}"}

@tool
def summarize_text(text: str, max_words: int = 400) -> dict:
    """Create a short extractive summary from long text (no LLM cost).
    Persists last_summary.
    """
    import textwrap
    words_left = max_words
    out: List[str] = []
    for para in re.split(r"\n{2,}", text or ""):
        para = para.strip()
        if not para:
            continue
        sent = re.split(r"(?<=[.!?])\s+", para)[0]
        w = sent.split()
        if not w:
            continue
        if len(w) > words_left:
            out.append(" ".join(w[:words_left]) + "…")
            words_left = 0
            break
        out.append(sent)
        words_left -= len(w)
        if words_left <= 0:
            break
    summary = "\n".join(textwrap.wrap(" ".join(out), width=100))

    st = _load_state()
    st.last_summary = summary
    _save_state(st)

    return {"status": "ok", "summary": summary}

def _ollama_json_outline(topic: Optional[str], notes: Optional[str], slide_target: int, model: str) -> str:
    """Internal helper: call ChatOllama to produce a strict-JSON slide outline; returns JSON string."""
    llm = ChatOllama(model=model, temperature=0.2)
    sys = f"""You output STRICT JSON for slide outlines only.
Schema:
{{
  "deck_title": string,
  "slides": [{{"title": string, "bullets": [string, ...]}}, ...]
}}
Rules:
- Exactly {slide_target} slides (1 title + {slide_target-1} content).
- Bullets are ≤12 words and non-redundant.
- Output JSON only. No code fences, no prose outside JSON.
"""
    parts = ["Create a deck outline."]
    if topic:
        parts.append(f"Topic: {topic}")
    if notes:
        parts.append(f"Use these notes:\n{(notes or '')[:8000]}")
    prompt = "\n".join(parts)
    resp = llm.invoke([SystemMessage(content=sys), HumanMessage(content=prompt)])
    content = (resp.content or "").strip()
    content = re.sub(r"^```json\s*|\s*```$", "", content, flags=re.S)
    try:
        json.loads(content)
    except Exception:
        content = json.dumps({
            "deck_title": topic or "Presentation",
            "slides": [{"title": topic or "Overview", "bullets": ["Outline generation failed; please retry."]}]
        }, ensure_ascii=False, indent=2)
    return content

_OUTLINE_CACHE: Dict[str, str] = {}

@tool
def slide_outline_json(
    topic: Optional[str] = None,
    notes: Optional[str] = None,
    slide_target: int = DEFAULTS["DEFAULT_SLIDES"],
    model = DEFAULTS["MODEL"]
) -> dict:
    """Create a JSON outline for a PPT (deck_title + slides[title, bullets]) using the LLM.
    Accepts optional notes and persists last_outline/topic. Always uses the local Ollama model.
    """
    # --- Force local Ollama model; ignore any external 'model' values (e.g., 'gpt-4') ---
    model = DEFAULTS["MODEL"]

    # --- Debounce key WITHOUT model so retries with different 'model' don't bypass cache ---
    key = _sig({
        "k": "outline",
        "t": topic or "",
        "n": hashlib.sha1((notes or "").encode()).hexdigest(),
        "s": slide_target,
    })
    if _recent_once(key, ttl=90) and key in _OUTLINE_CACHE:
        content = _OUTLINE_CACHE[key]
        st = _load_state()
        st.last_outline = content
        if topic:
            try:
                data = json.loads(content)
                st.topic = data.get("deck_title") or topic
            except Exception:
                st.topic = topic
        _save_state(st)
        return {"status": "cached", "outline_json": content}

    # --- Call the LLM to produce the outline (strict JSON) ---
    content = _ollama_json_outline(topic, notes, slide_target, model)
    _OUTLINE_CACHE[key] = content

    st = _load_state()
    st.last_outline = content
    if topic:
        try:
            data = json.loads(content)
            st.topic = data.get("deck_title") or topic
        except Exception:
            st.topic = topic
    _save_state(st)

    return {"status": "ok", "outline_json": content}


@tool
def build_ppt(outline_json: Union[str, dict], output_dir: Optional[str] = None) -> dict:
    """Render a PPTX file from a slide outline (string JSON or dict).

    Args:
        outline_json: The slide outline, either a JSON string or a parsed dict.
                      It must contain:
                        {
                          "deck_title": str,
                          "slides": [{"title": str, "bullets": [str, ...]}, ...]
                        }
        output_dir:   Optional directory to write the PPTX into. Defaults to DEFAULTS["OUTPUT_DIR"].

    Returns:
        dict with:
          - status: "ok" | "error" | "skipped"
          - path:   absolute PPTX filepath when status == "ok" (or when skipped due to debounce)
          - error:  error message when status == "error"
          - reason: reason when status == "skipped"
    """
    Presentation = DEPS["Presentation"]
    if Presentation is None:
        return {"status": "error", "error": "python-pptx not installed. pip install python-pptx"}

    # Accept both str and dict, normalize to string for hashing + parse to dict later.
    try:
        if isinstance(outline_json, dict):
            oj_str = json.dumps(outline_json, ensure_ascii=False)
        else:
            oj_str = str(outline_json)
    except Exception as e:
        return {"status": "error", "error": f"Could not serialize outline_json: {type(e).__name__}: {e}"}

    # Debounce on the normalized string to avoid duplicate builds
    sig = _sig({"k": "build", "outline": oj_str})
    if _recent_once(sig, ttl=120):
        st = _load_state()
        return {"status": "skipped", "reason": "recently_built", "path": st.ppt_path}

    # Parse into a dict the renderer can use
    try:
        data = json.loads(oj_str)
    except Exception:
        data = {"deck_title": "Presentation",
                "slides": [{"title": "Overview", "bullets": ["(Invalid JSON)"]}]}

    os.makedirs(output_dir or DEFAULTS["OUTPUT_DIR"], exist_ok=True)
    safe_title = re.sub(r"[^a-zA-Z0-9._-]+", "_", (data.get("deck_title") or "Presentation"))[:80] or "Presentation"
    fname = f"{safe_title}_{int(time.time())}.pptx"
    path = os.path.join(output_dir or DEFAULTS["OUTPUT_DIR"], fname)

    prs = Presentation()
    title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[1]
    bullet_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    slides = data.get("slides") or []
    for idx, s in enumerate(slides):
        layout = title_layout if idx == 0 else bullet_layout
        slide = prs.slides.add_slide(layout)
        shapes = slide.shapes
        title = shapes.title if hasattr(shapes, "title") else None
        if title:
            title.text = s.get("title") or f"Slide {idx+1}"
        # find a body placeholder for bullets
        body = None
        for shp in shapes:
            if shp.is_placeholder and getattr(shp, "placeholder_format", None) and shp.placeholder_format.type != 1:
                body = shp.text_frame
                break
        bullets = s.get("bullets") or []
        if body:
            body.clear()
            if bullets:
                body.text = bullets[0]
                for b in bullets[1:]:
                    p = body.add_paragraph(); p.text = b; p.level = 0

    from pathlib import Path as _P
    prs.save(path)
    st = _load_state()
    st.ppt_path = str(_P(path).resolve())
    _save_state(st)
    return {"status": "ok", "path": st.ppt_path}

@tool
def describe_latest_deck() -> dict:
    """Describe the most recently built deck without rebuilding it.

    Returns:
        dict: {
            "status": "ok" | "no_deck",
            "path": str | None,         # absolute PPTX path if available
            "topic": str | None,        # last known topic
            "headings": list[str]       # first ~8 slide titles if outline exists
        }
    """
    st = _load_state()
    if not st.ppt_path:
        return {"status": "no_deck"}
    headings: List[str] = []
    if st.last_outline:
        try:
            data = json.loads(st.last_outline)
            for s in (data.get("slides") or [])[:8]:
                headings.append(s.get("title") or s.get("heading") or "Slide")
        except Exception:
            pass
    return {"status": "ok", "path": st.ppt_path, "topic": st.topic, "headings": headings}

TOOLS = [read_pdf, summarize_text, slide_outline_json, build_ppt, describe_latest_deck]

# ---------------- Utility helpers ----------------
def ensure_dirs():
    os.makedirs(DEFAULTS["OUTPUT_DIR"], exist_ok=True)

def _pretty(obj: Any, max_len: int = 240) -> str:
    try:
        s = json.dumps(obj, ensure_ascii=False)
    except Exception:
        s = str(obj)
    return (s[:max_len] + "…") if len(s) > max_len else s

def _user_wants_build(txt: str) -> bool:
    t = (txt or "").lower()
    return (
        ("ppt" in t or "slides" in t or "presentation" in t)
        and ("build" in t or "make" in t or "create" in t)
    ) or ("based on this paper" in t or "use that" in t)

def _topic_from_text(txt: str) -> str:
    m = re.search(r"(about|on|regarding)\s+(.+)$", txt or "", flags=re.I)
    return (m.group(2).strip(" .!?\"'") if m else "Presentation")

# ---------------- Main / Agent loop ----------------
def main():
    parser = argparse.ArgumentParser(description="LangGraph ReAct agent (Ollama)")
    parser.add_argument("--model", default=DEFAULTS["MODEL"], help="Ollama model, e.g. gpt-oss:20b")
    parser.add_argument("--once", action="store_true", help="Read a single prompt from stdin then exit")
    parser.add_argument("--recursion", type=int, default=12, help="LangGraph recursion_limit per turn (default 12)")
    parser.add_argument("--verbose", action="store_true", help="Print step-by-step agent trace")
    parser.add_argument("--slides", type=int, default=DEFAULTS["DEFAULT_SLIDES"], help="Default slide count if user doesn't specify")
    args = parser.parse_args()

    # Reflect default slides into the prompt dynamically
    global SYSTEM_PROMPT
    SYSTEM_PROMPT = SYSTEM_PROMPT.replace("DEFAULT_SLIDES = "+str(DEFAULTS["DEFAULT_SLIDES"]), f"DEFAULT_SLIDES = {args.slides}")

    ensure_dirs()

    llm = ChatOllama(model=args.model, temperature=0.2)
    graph = create_react_agent(llm, tools=TOOLS)

    # Maintain history across turns; inject SystemMessage once
    history: List[Any] = [SystemMessage(content=SYSTEM_PROMPT)]

    def run_turn(user_text: str) -> str:
        nonlocal history
        history += [HumanMessage(content=user_text)]
        final = None
        step = 1
        # stream with recursion guard
        for event in graph.stream({"messages": history}, stream_mode="values", config={"recursion_limit": args.recursion}):
            final = event
            if not args.verbose:
                continue
            msgs = final.get("messages", [])
            if not msgs:
                continue
            last = msgs[-1]
            if isinstance(last, AIMessage):
                if getattr(last, "tool_calls", None):
                    for tc in last.tool_calls:
                        name = getattr(tc, "name", "tool")
                        arguments = getattr(tc, "args", {})
                        print(f"[step {step}] agent → tool: {name} args={_pretty(arguments)}")
                        step += 1
                else:
                    text = (last.content or "").strip()
                    if text:
                        print(f"[step {step}] agent message: {_pretty(text, 400)}")
                        step += 1
            elif isinstance(last, ToolMessage):
                name = getattr(last, "name", "tool")
                try:
                    content = json.loads(last.content) if isinstance(last.content, str) else last.content
                except Exception:
                    content = last.content
                brief = {}
                if isinstance(content, dict):
                    for k in ("status", "path", "topic", "headings", "reason", "outline_json"):
                        if k in content:
                            brief[k] = content[k] if k != "outline_json" else "(json outline)"
                else:
                    brief = {"content": content}
                print(f"[step {step}] tool → agent: {name} result={_pretty(brief)}")
                step += 1

        # After streaming, update history
        if final and "messages" in final and final["messages"]:
            history = final["messages"]
            if not history or not isinstance(history[0], SystemMessage):
                history = [SystemMessage(content=SYSTEM_PROMPT)] + history

            # If the model already produced content, return it
            last_ai = next((m for m in reversed(history) if isinstance(m, AIMessage)), None)
            if last_ai and ((getattr(last_ai, "tool_calls", None)) or (last_ai.content or "").strip()):
                return (last_ai.content or "").strip()

            # Fallback: user clearly wants a build but model didn’t call tools
            if _user_wants_build(user_text):
                st = _load_state()
                notes = st.last_summary or st.last_text
                topic = _topic_from_text(user_text)
                if notes:
                    outline_res = slide_outline_json.invoke({
                        "topic": topic,
                        "notes": notes[:8000],
                        "slide_target": args.slides,
                    })
                    if args.verbose:
                        print(f"[fallback] agent → tool: slide_outline_json args={_pretty({'topic': topic, 'slide_target': args.slides})}")
                        print(f"[fallback] tool → agent: slide_outline_json result={_pretty({k:('(json outline)' if k=='outline_json' else v) for k,v in outline_res.items()})}")
                    build_res = build_ppt.invoke({"outline_json": outline_res.get("outline_json", "{}")})
                    if args.verbose:
                        print(f"[fallback] agent → tool: build_ppt args={{'outline_json': '(json outline)'}}")
                        print(f"[fallback] tool → agent: build_ppt result={_pretty(build_res)}")
                    path = build_res.get("path")
                    return f"I created your deck about “{topic}”. File: {path}" if path else "I tried to build the deck but didn’t get a file path."
                else:
                    return "I can build the PPT. Provide a PDF/text file or a topic (e.g., “build a ppt about cats”)."

        # No final content
        return "(no response)"

    if args.once:
        user_text = sys.stdin.read().strip()
        reply = run_turn(user_text)
        print(reply)
        return

    print("ppt_agent (LangGraph ReAct + Ollama) — type your request. Ctrl+C to exit.")
    print(f"Model: {args.model}")
    while True:
        try:
            user_text = input("user > ").strip()
            if not user_text:
                continue
            reply = run_turn(user_text)
            print(f"agent> {reply}")
        except KeyboardInterrupt:
            print("\nBye.")
            break

if __name__ == "__main__":
    main()
