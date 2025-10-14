from __future__ import annotations
import os, json, re
from typing import TypedDict, Optional, List, Dict, Any, Literal

# --- LangChain / LangGraph ---
from langchain_ollama import ChatOllama
from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from langchain_core.runnables import RunnableLambda
from langgraph.graph import StateGraph, END

# --- Tools ---
from pydantic import BaseModel, Field
from langchain.tools import StructuredTool
import fitz  # PyMuPDF
from pptx import Presentation
import pathlib
import math

GENERIC_TITLE_PAT = re.compile(r"^(topic|section|slide)\s*\d+\b", re.I)

# =============================================================
# 1) Tool Implementations (pure Python + small LLM helpers)
# =============================================================

class ReadPDFIn(BaseModel):
    path: str = Field(..., description="Local path to a PDF file")
    pages: Optional[str] = Field(
        default=None,
        description='Optional page ranges like "1-3,5"; omit to read all.'
    )

def read_pdf(path: str, pages: Optional[str] = None) -> str:
    doc = fitz.open(path)
    selected: List[int] = []
    if pages:
        for part in pages.split(','):
            part = part.strip()
            if '-' in part:
                a, b = part.split('-')
                selected.extend(range(max(1, int(a)) - 1, int(b)))
            else:
                selected.append(max(1, int(part)) - 1)
        selected = [p for p in selected if 0 <= p < len(doc)]
    else:
        selected = list(range(len(doc)))
    chunks = []
    for i in selected:
        chunks.append(doc[i].get_text())
        # soft safety cap
        if sum(len(c) for c in chunks) > 25000:
            break
    text = "\n".join(chunks)
    # extra trim for token budget
    return text[:30000]

ReadPDF = StructuredTool.from_function(
    func=read_pdf,
    name="read_pdf",
    description=(
        "Read text from a research PDF. Use when the user provided a PDF path and you need its content."
    ),
    args_schema=ReadPDFIn,
)

# --- Summarize (uses the same Ollama LLM behind the scenes) ---
class SummarizeIn(BaseModel):
    text: str = Field(..., description="Raw content (PDF text or notes)")
    style: Literal["slides", "prose", "exec"] = Field(
        default="slides", description="Output style"
    )
    target_words: int = Field(400, description="Rough target length")

def summarize_text(text: str, style: str = "slides", target_words: int = 900) -> str:
    model = os.environ.get("OLLAMA_MODEL", "gpt-oss:20b")
    # small temperature helps the 20B model be less terse
    llm = ChatOllama(model=model, temperature=0.3)
    prompt = (
        "You are a subject-matter writer creating comprehensive notes for slide creation.\n"
        f"Topic: {text.strip()}\n\n"
        "Write ~900–1200 words of compact notes organized by clear sections using plain sentences "
        "(no markdown bullets). Include:\n"
        "- Background / definition\n"
        "- History / origins\n"
        "- Key concepts / taxonomy\n"
        "- Behavior / mechanisms / how it works\n"
        "- Real examples / case studies\n"
        "- Applications / impact\n"
        "- Health / risks / limitations\n"
        "- Best practices / tips\n"
        "- Recent trends / stats if known\n"
        "- Conclusion / key takeaways\n\n"
        "Rules:\n"
        "- No leading symbols like *, -, •\n"
        "- Short lines (1–2 sentences each), but dense with facts\n"
        "- Use concrete terms and names\n"
        "- Return ONLY the notes text (no JSON, no headings markup)\n"
        "- Do NOT prefix lines with section labels (e.g., 'Background/Definition:', 'Key Concepts/Taxonomy:', etc.). Just write the content."
    )
    return llm.invoke(prompt).content.strip()


Summarize = StructuredTool.from_function(
    func=summarize_text,
    name="summarize_text",
    description=(
        "Summarize content into concise bullets or prose. Use for 'summary only' requests or before making slides."
    ),
    args_schema=SummarizeIn,
)

# --- Slide outline JSON (use an LLM to structure nicely) ---
class OutlineIn(BaseModel):
    notes: str = Field(..., description="Summary notes to turn into slides")
    slide_target: int = Field(12, description="Approximate number of slides")

def densify_slide_with_llm(title: str, topic: str, notes: str, max_points: int = 5) -> tuple[str, list[str]]:
    model = os.environ.get("OLLAMA_MODEL", "gpt-oss:20b")
    llm = ChatOllama(model=model, temperature=0.3)
    prompt = (
        "Improve ONE presentation slide.\n"
        "Return STRICT JSON ONLY:\n"
        '{ "title": "Specific, informative title", "details": ["point 1", "point 2", "..."] }\n'
        "Rules:\n"
        "- 3 to 5 details; short phrases; no markdown bullets (*, -, •), no numbering.\n"
        "- Do NOT use phrases like 'Key takeaways include:' or 'Conclusion/Key Takeaways'.\n"
        "- Title must be informative (avoid 'Topic/Section/Slide N').\n\n"
        f"GlobalTopic: {topic}\n"
        f"SlideTitle: {title}\n\n"
        "Relevant notes:\n" + notes[:3500]
    )
    raw = llm.invoke(prompt).content.strip()
    m = re.search(r"\{[\s\S]*\}$", raw)
    js = m.group(0) if m else raw
    try:
        data = json.loads(js)
        t = re.sub(r"^[\s\*\-•·]+", "", str(data.get("title",""))).strip() or title
        pts = [re.sub(r"^[\s\*\-•·]+", "", str(x)).strip() for x in (data.get("details") or [])]
        pts = [p for p in pts if p][:max_points]
        if len(pts) < 3:
            raise ValueError("too few points")
        return t, pts
    except Exception:
        return title, synthesize_details_from_notes(title, notes, k=4)


def slide_outline_json(notes: str, slide_target: int = 12) -> Dict[str, Any]:
    model = os.environ.get("OLLAMA_MODEL", "gpt-oss:20b")
    llm = ChatOllama(model=model, temperature=0)

    schema_text = (
        '{ "deck_title": "string", '
        '"slides": [ { "title": "string", "subtitle": "string or empty", '
        '"details": ["short phrase", "..."] } ] }'
    )
    prompt = (
        "You are an expert presentation writer.\n"
        "Return JSON ONLY matching this schema:\n" + schema_text + "\n\n"
        "REQUIREMENTS:\n"
        f"- EXACTLY {slide_target} slides.\n"
        "- Slide 1: Welcome/Title slide with a strong deck_title (also use it as slide 1 title) and a short subtitle; details empty.\n"
        "- Slides 2..(N-2): Content slides with informative, specific titles (avoid 'Topic 1/2/...'); 3–5 concise details.\n"
        "- Slide (N-1): Conclusion (3–5 takeaways in details) with an informative title.\n"
        "- Slide N: Thanks for Watching (empty details).\n"
        "- Avoid markdown bullets (*, -, •). Use short, readable phrases.\n"
        "- Details must be standalone facts; do NOT include section labels.\n"
        "- Avoid incomplete sentences; end each detail with a period when appropriate.\n"
        "- JSON only. No extra text.\n\n"
        "NOTES:\n" + notes[:6000]
    )

    raw = llm.invoke(prompt).content.strip()
    m = re.search(r"\{[\s\S]*\}$", raw)
    json_text = m.group(0) if m else raw

    def _clean(s: str) -> str:
        return re.sub(r"^[\s\*\-•·]+", "", str(s or "")).strip()

    # Parse or start empty
    try:
        data = json.loads(json_text)
    except Exception:
        data = {"deck_title": make_deck_title_from_notes(notes), "slides": []}

    deck_title = _clean(data.get("deck_title") or "") or make_deck_title_from_notes(notes)
    # (bug fix) if you ever fallback to infer from user text, CALL the function
    if callable(deck_title):  # defensive in case of earlier mistakes
        deck_title = deck_title()

    slides = data.get("slides") or []
    sentences = split_sentences(notes)[:150]

    cleaned: List[Dict[str, Any]] = []
    for i, sl in enumerate(slides, start=1):
        title = _clean(sl.get("title", f"Slide {i}"))
        subtitle = _clean(sl.get("subtitle", ""))
        details = [_clean(x) for x in (sl.get("details") or []) if _clean(x)]

        # Special slides first
        if i == 1:
            if not deck_title:
                deck_title = title or "Presentation"
            title = deck_title
            details = []
        elif i == max(2, slide_target - 1):
            # Conclusion – ensure we have takeaways
            if not re.search(r"conclusion|takeaway|summary", title, re.I):
                title = "Conclusion"
            if len(details) < 3 or all("(add point)" in d for d in details):
                details = synthesize_details_from_notes("Key takeaways", notes, k=4)[:5]
        elif i == slide_target:
            title = "Thanks for Watching" if not re.search(r"thanks|thank", title, re.I) else title
            details = []
        else:
            # Regular slides – fix generic or thin content
            thin = (len(details) < 3) or all("(add point)" in d for d in details)
            if thin or is_generic_title(title):
                # (bug fix) pass topic correctly to densifier
                better_title, better_details = densify_slide_with_llm(title, deck_title, notes, max_points=5)
                if not is_generic_title(better_title):
                    title = better_title
                details = [d for d in better_details if d]

            if (len(details) < 3) or all("(add point)" in d for d in details):
                # Deterministic backup: synthesize from notes
                details = synthesize_details_from_notes(title, notes, k=4)

            # Final clamp: 3–5 short details
            details = [re.sub(r"^[\s\*\-•·]+", "", d).strip() for d in details if d.strip()]
            while len(details) < 3:
                details.append("(add point)")
            details = details[:5]

        # (bug fix) append exactly once per slide (no double-append)
        cleaned.append({"title": title, "subtitle": subtitle, "details": details})

    # Pad/truncate to exactly N (preserve first + last two)
    if len(cleaned) < slide_target:
        want_mid = max(0, slide_target - 3)
        suggested = keyphrase_titles_from_notes(notes, want_mid)
        insert_at = max(1, len(cleaned) - 1)
        while len(cleaned) < slide_target:
            idx = len(cleaned) - 1
            t = suggested[idx % max(1, len(suggested))] if suggested else f"Key Topic {idx+1}"
            cleaned.insert(insert_at, {"title": t, "subtitle": "", "details": synthesize_details_from_notes(t, notes, k=4)})
    elif len(cleaned) > slide_target:
        first = cleaned[0]
        tail = cleaned[-2:] if slide_target >= 3 else []
        mid = cleaned[1:-2][:max(0, slide_target - (1 + len(tail)))]
        cleaned = [first] + mid + tail

    # Final guarantees (Welcome/Conclusion/Thanks)
    if slide_target >= 1:
        cleaned[0]["title"] = deck_title
        cleaned[0]["subtitle"] = cleaned[0].get("subtitle") or f"An overview of {deck_title}"
        cleaned[0]["details"] = []
    if slide_target >= 3:
        cleaned[-2]["title"] = "Conclusion"
        if len(cleaned[-2].get("details", [])) < 3:
            cleaned[-2]["details"] = synthesize_details_from_notes("Key takeaways", notes, k=4)[:5]
        cleaned[-1]["title"] = "Thanks for Watching"
        cleaned[-1]["details"] = []

    cleaned = polish_outline(cleaned, notes)
    cleaned = allocate_details_across_slides(cleaned, notes)
    deck_title = deck_title or (cleaned[1]["title"] if len(cleaned) > 1 else "Presentation")
    return {"deck_title": deck_title, "slides": cleaned}

# --- Build PPT ---
class BuildIn(BaseModel):
    outline_json: Dict[str, Any] = Field(..., description="Slide outline JSON")
    path: str = Field(default="output.pptx", description="Output path for .pptx")

def build_ppt(outline_json: Dict[str, Any], path: str = "output.pptx") -> str:
    prs = Presentation()

    def clean_line(s: str) -> str:
        return re.sub(r"^[\s\*\-•·]+", "", s or "").strip()

    slides = outline_json.get("slides", []) or []
    deck_title = clean_line(outline_json.get("deck_title", ""))

    if not slides:
        slides = [{"title": deck_title or "Welcome", "subtitle": "", "details": []},
                  {"title": "Conclusion", "subtitle": "", "details": ["(add point)", "(add point)", "(add point)"]},
                  {"title": "Thanks for Watching", "subtitle": "", "details": []}]

    # Prefer deck_title for filename; fallback to slide 1/2
    name_source = deck_title or clean_line(slides[0].get("title", "")) or (clean_line(slides[1].get("title", "")) if len(slides) > 1 else "Presentation")
    safe_title = re.sub(r"[^\w\-]+", "_", name_source).strip("_") or "Presentation"
    if not path or path == "output.pptx":
        path = f"{safe_title}.pptx"
        base, ext = os.path.splitext(path)
        i = 2
        while os.path.exists(path):
            path = f"{base}_{i}{ext}"
            i += 1

    total = len(slides)

    for idx, sl in enumerate(slides, start=1):
        title = clean_line(sl.get("title", f"Slide {idx}"))
        subtitle = clean_line(sl.get("subtitle", ""))
        details = [clean_line(x) for x in (sl.get("details") or []) if clean_line(x)]

        is_welcome = (idx == 1)
        is_thanks  = (idx == total)
        is_concl   = (idx == max(2, total-1))

        if is_welcome:
            layout = prs.slide_layouts[0]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = title
            if len(s.placeholders) > 1:
                s.placeholders[1].text = subtitle
            else:
                box = s.shapes.add_textbox(
                    left=prs.slide_width//10,
                    top=prs.slide_height//2,
                    width=prs.slide_width*8//10,
                    height=prs.slide_height//5
                )
                box.text_frame.text = subtitle or ""
            continue

        layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = title
        body = s.placeholders[1]

        if is_thanks:
            # Thanks slide should have no bullets
            body.text = subtitle or ""
            continue

        # Non-thanks slides: enforce at least 3 bullets on regular and conclusion
        if not details:
            details = ["(add point)", "(add point)", "(add point)"]
        elif (not is_welcome) and (not is_thanks) and len(details) < 3:
            details += ["(add point)"] * (3 - len(details))

        tf = body.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.text = details[0]
        for d in details[1:]:
            p = tf.add_paragraph()
            p.text = d

    prs.save(path)
    return os.path.abspath(path)


Build = StructuredTool.from_function(
    func=build_ppt,
    name="build_ppt",
    description=(
        "Render a PowerPoint (.pptx) from slide_outline_json. Final step only."
    ),
    args_schema=BuildIn,
)
class OutlineIn(BaseModel):
    notes: str
    slide_target: int = 10

Outline = StructuredTool.from_function(
    func=slide_outline_json,
    name="slide_outline_json",
    description="Generate slide outline JSON (title, subtitle, details) from summarized notes",
    args_schema=OutlineIn,
)
TOOLS: Dict[str, StructuredTool] = {
    "read_pdf": ReadPDF,
    "summarize_text": Summarize,
    "slide_outline_json": Outline,
    "build_ppt": Build,
}

# =============================================================
# 2) LangGraph State & Nodes
# =============================================================

class AgentState(TypedDict, total=False):
    messages: List[Dict[str, str]]         # chat history [{role, content}]
    pdf_text: Optional[str]
    notes: Optional[str]
    outline: Optional[Dict[str, Any]]
    ppt_path: Optional[str]
    next_action: Optional[str]
    next_args: Optional[Dict[str, Any]]
    auto: bool                             # if True, don't ask before building

# Helper to get last user message
def last_user_text(state: AgentState) -> str:
    for m in reversed(state.get("messages", [])):
        if m["role"] == "user":
            return m["content"]
    return ""

BAD_PHRASE_PAT = re.compile(
    r"(?i)\b(key takeaways include:|conclusion/key takeaways|conclusion:|^[-*•]+\s*$|\*\*|^\*|\s\*$)"
)
SECTION_PREFIX_PAT = re.compile(
    r"(?i)^\s*(background/?definition|history/?origins|key concepts/?taxonomy|"
    r"real examples/?case studies|applications/?impact|health/?risks/?limitations|"
    r"best practices/?tips|recent trends/?stats)\s*[:\-]*\s*"
)

def fix_truncated_line(line: str, notes: str) -> str:
    s = line.rstrip()
    # heuristic: ends with a 1–3 letter fragment
    if re.search(r"[A-Za-z]{1,3}$", s) and not s.endswith((".", "!", "?")):
        # find a note sentence that begins with this prefix (case-insensitive)
        prefix = re.sub(r"\s+", " ", s[-40:]).strip().lower()
        for sent in split_sentences(notes):
            low = sent.lower()
            if low.startswith(prefix):
                return sent[:120].rstrip()
        # otherwise trim the dangling fragment
        s = re.sub(r"\s+[A-Za-z]{1,3}$", "", s).rstrip()
    return s

def clean_detail_line(s: str) -> str:
    s = re.sub(r"^[\s\*\-•·]+", "", (s or "").strip())
    s = SECTION_PREFIX_PAT.sub("", s)
    s = BAD_PHRASE_PAT.sub("", s)
    s = re.sub(r"\*\*+", "", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s

def jaccard(a: str, b: str) -> float:
    wa, wb = set(re.findall(r"[a-z0-9]+", a.lower())), set(re.findall(r"[a-z0-9]+", b.lower()))
    if not wa or not wb: return 0.0
    return len(wa & wb) / len(wa | wb)

def compress_details(lines: list[str], min_keep: int = 3, max_keep: int = 5, sim_thresh: float = 0.6) -> list[str]:
    out = []
    for ln in (clean_detail_line(x) for x in lines):
        if not ln: continue
        if any(jaccard(ln, y) >= sim_thresh for y in out):  # drop near-duplicates
            continue
        out.append(ln)
        if len(out) >= max_keep: break
    while len(out) < min_keep:
        out.append("(add point)")
    return out[:max_keep]

def polish_outline(cleaned: list[Dict[str, Any]], notes: str) -> list[Dict[str, Any]]:
    if not cleaned: return cleaned
    n = len(cleaned)
    # Fix middle slides (2..n-2)
    for i in range(1, max(1, n-2)):
        sl = cleaned[i]
        sl["title"] = repair_title(sl.get("title", ""), notes)
        details = sl.get("details") or []
        details = compress_details(details, min_keep=3, max_keep=5)
        if all("(add point)" in d for d in details):
            # synthesize from notes if still empty-ish
            details = synthesize_details_from_notes(sl["title"], notes, k=4)
            details = compress_details(details, min_keep=3, max_keep=5)
        sl["details"] = details

    # Ensure Conclusion (n-1)
    if n >= 2:
        concl = cleaned[-2]
        concl["title"] = "Conclusion"
        concl["details"] = compress_details(
            concl.get("details") or synthesize_details_from_notes("Key takeaways", notes, k=4),
            min_keep=3, max_keep=5
        )

    # Ensure Thanks (n)
    cleaned[-1]["title"] = "Thanks for Watching"
    cleaned[-1]["details"] = []

    # Welcome (1): keep subtitle tidy
    cleaned[0]["subtitle"] = re.sub(r"\s+", " ", (cleaned[0].get("subtitle") or f"An overview of {cleaned[0].get('title','')}")).strip()
    cleaned[0]["details"] = []
    return cleaned

def repair_title(t: str, notes: str) -> str:
    t = (t or "").strip()
    if not t or re.fullmatch(r"(have been|dogs highly|highly social|other dogs|dogs also)", t, re.I) or len(t.split()) < 2:
        # propose a better title from notes’ keyphrases
        cands = keyphrase_titles_from_notes(notes, 1)
        return cands[0] if cands else "Key Topic"
    # title-case and strip junk
    t = re.sub(r"\s+", " ", t)
    return t[0].upper() + t[1:]

def extract_pdf_path(text: str) -> str | None:
    """Return a plausible PDF path from a user message (Windows-friendly)."""
    # Match quoted paths: "D:\dir\file.pdf"
    m = re.search(r'"([^"]+\.pdf)"', text, re.I)
    if m:
        return m.group(1)
    # Match unquoted Windows path or bare name ending with .pdf
    m = re.search(r'([A-Za-z]:\\[^\s]+\.pdf)', text, re.I)
    if m:
        return m.group(1)
    m = re.search(r'(\S+\.pdf)\b', text, re.I)
    if m:
        return m.group(1)
    return None

def wants_summary_only_now(state: AgentState) -> bool:
    return "summary only" in (last_user_text(state).lower())

def infer_topic_from_user_text(state: AgentState) -> str:
    t = last_user_text(state)
    t = re.sub(r"\b(make|create|build|slides?|ppt|presentation|about|on)\b", "", t, flags=re.I)
    t = re.sub(r"\s+", " ", t).strip()
    return t[:80].title() or "Presentation"

# Helper to append a message to state
def add_message(state: AgentState, role: str, content: str) -> AgentState:
    msgs = list(state.get("messages", []))
    msgs.append({"role": role, "content": content})
    state["messages"] = msgs
    return state

def infer_topic_from_user_text(state: AgentState) -> str:
    t = last_user_text(state)
    # strip common verbs
    t = re.sub(r"\b(make|create|build|slides?|ppt|presentation|about|on)\b", "", t, flags=re.I)
    t = re.sub(r"\s+", " ", t).strip()
    # Title Case
    return t[:80].title() or "Presentation"

def split_sentences(text: str) -> list[str]:
    # very light splitter; avoids bringing in external deps
    parts = re.split(r'(?<=[.!?])\s+', text.strip())
    return [p.strip() for p in parts if p.strip()]

def normalize_line_global(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "").strip().lower())

def compress_details_unique(lines: list[str], global_seen: set[str],
                            min_keep: int = 3, max_keep: int = 5, sim_thresh: float = 0.6) -> list[str]:
    out = []
    for ln in (clean_detail_line(x) for x in lines):
        if not ln:
            continue
        norm = normalize_line_global(ln)
        if norm in global_seen:
            continue
        if any(jaccard(ln, y) >= sim_thresh for y in out):
            continue
        out.append(ln)
        global_seen.add(norm)
        if len(out) >= max_keep:
            break
    while len(out) < min_keep:
        out.append("(add point)")
    return out[:max_keep]

def allocate_details_across_slides(cleaned: list[Dict[str, Any]], notes: str) -> list[Dict[str, Any]]:
    if not cleaned: return cleaned
    n = len(cleaned)
    global_seen: set[str] = set()

    # Reserve welcome & thanks; seed global_seen with their lines (usually none)
    for idx in [0, n-1]:
        for ln in cleaned[idx].get("details", []) or []:
            global_seen.add(normalize_line_global(ln))

    # Walk through content slides in order, compressing uniquely and filling from notes with exclude=global_seen
    for i in range(1, max(1, n-1)):
        if i == n-1:  # thanks
            break
        sl = cleaned[i]
        if sl.get("title", "").lower() == "conclusion":
            # handle conclusion last in this loop
            continue

        # 1) Unique-compress what's there
        uniq = compress_details_unique(sl.get("details") or [], global_seen, min_keep=0, max_keep=5)

        # 2) If <3, synthesize from notes excluding global_seen
        if len(uniq) < 3:
            extra = synthesize_details_from_notes(sl.get("title", ""), notes, k=5-len(uniq), exclude=global_seen)
            for e in extra:
                norm = normalize_line_global(e)
                if norm in global_seen:
                    continue
                uniq.append(e)
                global_seen.add(norm)
                if len(uniq) >= 5:
                    break

        # 3) Guarantee 3–5
        while len(uniq) < 3:
            uniq.append("(add point)")
        sl["details"] = uniq[:5]

    # Finally ensure conclusion has unique lines too
    if n >= 2:
        concl = cleaned[-2]
        uniq = compress_details_unique(concl.get("details") or [], global_seen, min_keep=0, max_keep=5)
        if len(uniq) < 3:
            extra = synthesize_details_from_notes("Key takeaways", notes, k=5-len(uniq), exclude=global_seen)
            for e in extra:
                norm = normalize_line_global(e)
                if norm in global_seen:
                    continue
                uniq.append(e)
                global_seen.add(norm)
                if len(uniq) >= 5:
                    break
        while len(uniq) < 3:
            uniq.append("(add point)")

        # NEW: if placeholders remain, try LLM densify, then fallback synthesize without exclude
        if all("(add point)" in d for d in uniq):
            try:
                bt, bd = densify_slide_with_llm(sl.get("title",""), cleaned[0].get("title","Presentation"), notes, max_points=5)
                cand = [clean_detail_line(fix_truncated_line(x, notes)) for x in bd]
                uniq = [x for x in cand if x][:5] or uniq
            except Exception:
                pass
            if all("(add point)" in d for d in uniq):
                extra = synthesize_details_from_notes(sl.get("title",""), notes, k=5, exclude=set())  # <— no exclude now
                cand = [clean_detail_line(fix_truncated_line(x, notes)) for x in extra]
                uniq = [x for x in cand if x][:5] or uniq

        sl["details"] = uniq[:5]
    return cleaned


def keyphrase_titles_from_notes(notes: str, k: int) -> list[str]:
    text = re.sub(r"[^a-z0-9\s]", " ", notes.lower())
    words = [w for w in text.split() if 2 <= len(w) <= 20]
    if not words:
        return []
    # stopwords-lite
    stop = set("a an the of to for and or in on with at by from as is are be was were this that these those it its their our your you we us".split())
    words = [w for w in words if w not in stop]
    # generate bi/tri-grams
    grams = []
    for n in (3, 2):
        for i in range(len(words)-n+1):
            grams.append(" ".join(words[i:i+n]))
    # count + rank
    freq: Dict[str,int] = {}
    for g in grams:
        freq[g] = freq.get(g, 0) + 1
    ranked = sorted(freq.items(), key=lambda x: (-x[1], -len(x[0])))
    titles = []
    for g, _ in ranked:
        # title case, de-dup by substring overlap
        cand = " ".join(w.capitalize() for w in g.split())
        if any(cand in t or t in cand for t in titles):
            continue
        titles.append(cand)
        if len(titles) >= k:
            break
    return titles

def is_generic_title(t: str) -> bool:
    return bool(GENERIC_TITLE_PAT.match((t or "").strip()))

def make_deck_title_from_notes(notes: str) -> str:
    # build a readable deck name from top keyphrases or first good sentence
    keys = keyphrase_titles_from_notes(notes, 1)
    if keys:
        return keys[0]
    sents = split_sentences(notes)
    if sents:
        return sents[0][:60]
    return "Presentation"

def tokenize(s: str) -> list[str]:
    return re.findall(r"[a-zA-Z0-9]+", (s or "").lower())

def synthesize_details_from_notes(title: str, notes: str, k: int = 4, exclude: Optional[set[str]] = None) -> list[str]:
    exclude = exclude or set()
    title_terms = set(w for w in tokenize(title) if len(w) > 2)
    sents = split_sentences(notes)
    if not sents:
        return ["(add point)"] * max(3, k)

    scored = []
    for s in sents:
        raw = s.strip()
        key = re.sub(r"\s+", " ", raw.lower())
        if key in exclude:
            continue
        ws = set(w for w in tokenize(raw) if len(w) > 2)
        overlap = len(ws & title_terms)
        length = len(raw)
        length_score = -abs(length - 90) / 90.0
        score = overlap * 3.0 + length_score
        scored.append((score, raw))

    scored.sort(reverse=True, key=lambda x: x[0])

    picked = []
    seen_local = set()
    for _, raw in scored:
        line = re.sub(r"^[\s\*\-•·]+", "", raw)
        norm = re.sub(r"\s+", " ", line.lower()).strip()
        if not norm or norm in seen_local or norm in exclude:
            continue
        # shorten nicely
        if len(line) > 100:
            cut = re.search(r"[;,:–—\-]\s", line[60:100])
            line = (line[:60 + cut.start()].strip() if cut else line[:90].rstrip())
        picked.append(line)
        seen_local.add(norm)
        if len(picked) >= k:
            break

    while len(picked) < max(3, k):
        picked.append("(add point)")
    return picked[: max(3, k)]


# --- Planner Node: LLM decides next action (which tool/chat/finish) ---
def planner_node(state: AgentState) -> AgentState:
    """
    Decide the next step:
      - read_pdf → summarize_text → slide_outline_json → build_ppt
      - handle "summary only", "regenerate outline", explicit slide counts
      - avoid summarize→summarize loops
    """
    model = os.environ.get("OLLAMA_MODEL", "gpt-oss:20b")
    llm = ChatOllama(model=model, temperature=0)

    # --- gather context
    user_text = last_user_text(state).strip()
    text_l = user_text.lower()

    # clear stale intent before we plan
    state["next_action"] = None
    state["next_args"] = {}

    # early exit: if PPT already built, finish this turn
    if state.get("ppt_path"):
        state["next_action"] = "finish"
        state["next_args"] = {}
        return state

    def extract_slide_count(s: str, default: int = 11) -> int:
        m = re.search(r"(\d+)\s*(slides?|ppt)", s.lower())
        return int(m.group(1)) if m else default

    # quick intent flags
    wants_ppt = any(k in text_l for k in [
        "ppt", "slides", "slide deck", "powerpoint", "power point",
        "presentation", "presentaion"
    ])
    wants_summary_only = ("summary only" in text_l) or (("summary" in text_l) and not wants_ppt)
    mentions_pdf = ".pdf" in text_l
    asks_regen_outline = "regenerate outline" in text_l or "re-generate outline" in text_l or "redo outline" in text_l
    asks_build_now = "build slides now" in text_l or text_l == "build slides" or "build now" in text_l

    slide_target = extract_slide_count(text_l, default=11)
    pdf_path = extract_pdf_path(user_text) if mentions_pdf else None

    # --- high-priority control commands
    if asks_regen_outline and state.get("notes"):
        state["next_action"] = "slide_outline_json"
        state["next_args"] = {"notes": state["notes"], "slide_target": slide_target}
        return state

    if asks_build_now and state.get("outline"):
        state["next_action"] = "build_ppt"
        state["next_args"] = {"outline_json": state["outline"], "path": "output.pptx"}
        return state

    # --- primary plan via LLM (lightweight) – JSON {action, args}
    tool_descriptions = "\n".join([f"- {name}: {tool.description}" for name, tool in TOOLS.items()])
    context_bits = []
    if state.get("pdf_text"): context_bits.append("pdf_text: present")
    if state.get("notes"):    context_bits.append("notes: present")
    if state.get("outline"):  context_bits.append("outline: present")
    context_line = ", ".join(context_bits) or "(no artifacts yet)"

    system = SystemMessage(content=(
        "You are a tool-using assistant. Decide the next best action.\n"
        "Allowed actions: read_pdf, summarize_text, slide_outline_json, build_ppt, chat, finish.\n"
        "If user said 'summary only', do not build slides.\n"
        "If building slides and no outline yet, create outline first.\n"
        "Return JSON ONLY: {\"action\": <one_of_allowed>, \"args\": {...}}"
    ))
    human = HumanMessage(content=(
        f"TOOLS:\n{tool_descriptions}\n\n"
        f"ARTIFACTS: {context_line}\n"
        f"User message:\n{user_text}\n"
        "JSON only, no prose."
    ))

    action, args = "chat", {}
    try:
        resp = llm.invoke([system, human]).content.strip()
        m = re.search(r"\{[\s\S]*\}$", resp)
        js = m.group(0) if m else resp
        data = json.loads(js)
        a = (data.get("action") or "chat").lower()
        if a in {"read_pdf","summarize_text","slide_outline_json","build_ppt","chat","finish"}:
            action = a
            args = data.get("args") or {}
    except Exception:
        # --- heuristic fallback (robust / deterministic)
        if wants_summary_only:
            if state.get("pdf_text"):
                action, args = "summarize_text", {"text": state["pdf_text"], "style": "slides", "target_words": 1100}
            elif pdf_path:
                action, args = "read_pdf", {"path": pdf_path}
            else:
                topic = infer_topic_from_user_text(state)
                topic_seed = f"{topic}\n\nCreate comprehensive notes about this topic for slides."
                action, args = "summarize_text", {"text": topic_seed, "style": "slides", "target_words": 1100}

        elif mentions_pdf and state.get("pdf_text"):
            action, args = "summarize_text", {"text": state["pdf_text"], "style": "slides", "target_words": 1100}

        elif mentions_pdf and pdf_path:
            action, args = "read_pdf", {"path": pdf_path}

        elif wants_ppt:
            if state.get("notes") and not state.get("outline"):
                action, args = "slide_outline_json", {"notes": state["notes"], "slide_target": slide_target}
            elif state.get("outline"):
                action, args = "build_ppt", {"outline_json": state["outline"], "path": "output.pptx"}
            else:
                topic = infer_topic_from_user_text(state)
                topic_seed = f"{topic}\n\nCreate comprehensive notes about this topic for slides."
                action, args = "summarize_text", {"text": topic_seed, "style": "slides", "target_words": 1100}

        else:
            action, args = "chat", {}

    # -------- COERCION BLOCK (prevents loops, advances pipeline) --------
    if wants_ppt and not wants_summary_only:
        # If we already have notes but the plan chose summarize again → make outline
        if state.get("notes") and not state.get("outline") and action == "summarize_text":
            action = "slide_outline_json"
            args = {"notes": state["notes"], "slide_target": slide_target}

        # If we have outline → build now (even if model said outline/summarize)
        if state.get("outline") and action in {"summarize_text", "slide_outline_json"}:
            action = "build_ppt"
            args = {"outline_json": state["outline"], "path": "output.pptx"}

    # If we’re not making PPT and already have notes, avoid re-summarizing forever
    if action == "summarize_text" and state.get("notes") and not wants_summary_only and not wants_ppt:
        action, args = "chat", {}

    state["next_action"] = action
    state["next_args"] = args
    return state


# --- Tool Executor Node: run the chosen tool (if any) ---
def tool_node(state: AgentState) -> AgentState:
    action = state.get("next_action")
    args = dict(state.get("next_args") or {})  # copy to mutate safely

    # --- Pre-fill missing args defensively ---
    if action == "summarize_text":
        if not args.get("text"):
            # Prefer already-read PDF text; otherwise seed from user topic
            seed = state.get("pdf_text") or f"Create slide-ready notes about: {last_user_text(state)}"
            args["text"] = seed
        args.setdefault("style", "slides")
        args.setdefault("target_words", 900)

    elif action == "slide_outline_json":
        if not args.get("notes"):
            # Prefer existing notes; otherwise fall back to the user's last message
            seed_notes = state.get("notes") or f"Key points about: {last_user_text(state)}"
            args["notes"] = seed_notes
        # Keep any slide_target the planner set; default to 10
        args.setdefault("slide_target", 10)

    elif action == "build_ppt":
        if not args.get("outline_json"):
            if state.get("outline"):
                args["outline_json"] = state["outline"]
                args.setdefault("path", "output.pptx")
            else:
                # Can't build without an outline — give a friendly message and end the turn
                state = add_message(
                    state, "assistant",
                    "I need a slide outline before building. Say 'outline it' or 'build slides' after I prepare the outline."
                )
                state["next_action"] = "finish"
                state["next_args"] = {}
                return state

    # --- Execute the (now valid) tool call ---
    if action in TOOLS:
        tool = TOOLS[action]
        result = tool.invoke(args)

        if action == "read_pdf":
            state["pdf_text"] = result
            state = add_message(state, "assistant", "Read PDF ✔ (text cached).")

        elif action == "summarize_text":
            state["notes"] = result
            shown = (result[:2000] + "...") if len(result) > 2000 else result
            state = add_message(
                state, "assistant",
                f"Summary ready ✔:\n\n{shown}\n\n"
                "Say 'build slides' to turn this into a deck, or 'summary only' to end."
            )
            if wants_summary_only_now(state):
                state["next_action"] = "finish"
                state["next_args"] = {}
                return state

        elif action == "slide_outline_json":
            state["outline"] = result

            # pretty-print JSON to chat
            try:
                pretty = json.dumps(result, indent=2, ensure_ascii=False)
            except Exception:
                pretty = str(result)

            # also save to a file for easy inspection
            try:
                with open("outline_debug.json", "w", encoding="utf-8") as f:
                    f.write(pretty)
                saved_note = " (saved to outline_debug.json)"
            except Exception:
                saved_note = ""

            state = add_message(
                state,
                "assistant",
                "Slide outline JSON prepared ✔" + saved_note + ". Here's the JSON:\n\n"
                f"```json\n{pretty}\n```\n"
                "Say **'build slides'** to render this outline, or **'regenerate outline'** to try again."
            )

        elif action == "build_ppt":
            state["ppt_path"] = result if isinstance(result, str) else args.get("path", "output.pptx")
            state = add_message(state, "assistant", f"PPT built ✔ -> {state['ppt_path']}")
            # Make this turn terminal to avoid tool→plan→build loops
            state["next_action"] = "finish"
            state["next_args"] = {}
            return state

    # Clear to avoid stale directives on the next plan step
    state["next_action"] = None
    state["next_args"] = {}
    return state

# --- Chat Node: simple chat response (no tools) ---
def chat_node(state: AgentState) -> AgentState:
    """Plain conversational reply (no tools). Ends the turn via chat→END edge."""
    model = os.environ.get("OLLAMA_MODEL", "gpt-oss:20b")
    llm = ChatOllama(model=model, temperature=0)

    # Short system guidance to keep answers crisp
    sys = SystemMessage(content=(
        "You are a helpful research→slides assistant. "
        "Answer briefly and clearly. If the user asks for PPTs or summaries, "
        "you may propose a plan, but do not call tools from this node."
    ))

    # Include a small window of history (last few turns)
    history = [sys]
    for m in state.get("messages", [])[-6:]:
        if m["role"] == "user":
            history.append(HumanMessage(content=m["content"]))
        else:
            history.append(AIMessage(content=m["content"]))

    reply = llm.invoke(history).content
    state = add_message(state, "assistant", reply)

    # Clear any stale next_* so the router doesn't reuse them
    state["next_action"] = None
    state["next_args"] = {}
    return state

# --- Decider for graph edges ---

def choose_route(state: AgentState) -> Literal["tool", "chat", "finish"]:
    action = (state.get("next_action") or "chat").lower()
    if action == "finish":
        return "finish"
    if action in TOOLS:
        return "tool"
    if action == "chat":
        return "chat"
    return "chat"


# =============================================================
# 3) Build Graph Application
# =============================================================

def build_app():
    g = StateGraph(AgentState)
    g.add_node("plan", planner_node)
    g.add_node("tool", tool_node)
    g.add_node("chat", chat_node)

    g.set_entry_point("plan")

    g.add_conditional_edges("plan", choose_route, {
        "tool": "tool",
        "chat": "chat",
        "finish": END,
    })
    g.add_edge("tool", "plan")
    g.add_edge("chat", END)

    app = g.compile()
    return app

# =============================================================
# 4) Simple CLI loop
# =============================================================

if __name__ == "__main__":
    model = os.environ.get("OLLAMA_MODEL", "gpt-oss:20b")
    print(f"🧠 Using Ollama model: {model}")
    app = build_app()

    state: AgentState = {"messages": []}

    print("\n📽️ LangGraph PPT Agent ready. Type your request (or 'exit').")
    print("Examples:\n - summary only for D:/papers/attention.pdf\n - make slides about diffusion models (10 slides)\n - read D:/papers/paper.pdf pages 1-3 and summarize\n - build slides now\n")

    steps = 0
    while True:
        try:
            user = input("You> ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\nBye!")
            break
        if user.lower() in {"exit", "quit"}:
            print("Bye!")
            break
        state = add_message(state, "user", user)
        state = app.invoke(state)   # exactly once
        if state.get("messages"):
            last = state["messages"][-1]
            if last["role"] == "assistant":
                print("Assistant>", last["content"])
