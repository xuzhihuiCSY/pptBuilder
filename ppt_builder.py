import sys
import os
import json
import re
from typing import List, Optional, Any, Tuple
import datetime
import requests
from pydantic import BaseModel, Field, ValidationError
from pypdf import PdfReader

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import PydanticOutputParser
from langchain_text_splitters import RecursiveCharacterTextSplitter

from langchain_ollama import ChatOllama


# =========================
# Schemas
# =========================
class Slide(BaseModel):
    title: str
    subtitle: str = Field("", description="<= 12 words")
    details: List[str] = Field(default_factory=list, description="3–6 concise bullets")
    speaker_notes: Optional[str] = None


class Deck(BaseModel):
    topic: str
    audience: str
    tone: str
    slide_count: int
    slides: List[Slide]


# =========================
# Helpers
# =========================
def ensure_ollama_running(base_url: str = "http://localhost:11434"):
    try:
        r = requests.get(f"{base_url}/api/tags", timeout=2)
        if r.status_code != 200:
            raise RuntimeError(f"Ollama responded with status {r.status_code}.")
    except Exception as e:
        raise RuntimeError(
            "❌ Couldn't reach Ollama at http://localhost:11434.\n"
            "Start it in another terminal:\n"
            "    ollama serve\n"
        ) from e


def extract_json_block(text: str) -> str:
    """Extract the first JSON object from text (in case the model adds prose)."""
    start = text.find("{"); 
    end = text.rfind("}")
    if start != -1 and end != -1 and end > start:
        return text[start:end + 1]
    return text


def coerce_bullets(details: Any) -> List[str]:
    """Ensure details is a list[str]. Flatten dict items into strings if needed."""
    out: List[str] = []
    if not details:
        return out
    if isinstance(details, str):
        return [details]
    if isinstance(details, list):
        for item in details:
            if isinstance(item, str):
                s = item.strip()
                if s:
                    out.append(s)
            elif isinstance(item, dict):
                parts = []
                for k in ("value", "myth", "fact", "title", "key", "point", "text"):
                    v = item.get(k)
                    if isinstance(v, str) and v.strip():
                        parts.append(v.strip())
                for k in ("description", "explanation", "detail"):
                    v = item.get(k)
                    if isinstance(v, str) and v.strip():
                        parts.append(v.strip())
                if parts:
                    out.append(" — ".join(parts))
            else:
                out.append(str(item))
    else:
        out.append(str(details))
    return [b for b in (s.strip() for s in out) if b]


def repair_deck_obj(obj: Any) -> Any:
    """Normalize raw JSON to Deck schema expectations."""
    if not isinstance(obj, dict):
        return obj
    slides = obj.get("slides", [])
    fixed_slides = []
    for s in slides:
        if not isinstance(s, dict):
            continue
        title = (s.get("title") or "").strip() or "Untitled"
        subtitle = (s.get("subtitle") or "").strip() or "Overview"
        details = coerce_bullets(s.get("details", []))
        while len(details) < 3:
            details.append("Key point")
        details = details[:6]
        speaker_notes = s.get("speaker_notes") or None
        fixed_slides.append(
            {"title": title, "subtitle": subtitle, "details": details, "speaker_notes": speaker_notes}
        )
    topic = str(obj.get("topic") or "Presentation").strip()
    audience = str(obj.get("audience") or "general").strip()
    tone = str(obj.get("tone") or "informative").strip()
    slide_count = int(obj.get("slide_count") or len(fixed_slides) or 8)
    return {"topic": topic, "audience": audience, "tone": tone, "slide_count": slide_count, "slides": fixed_slides}


def autodetect_pdf(path: Optional[str]) -> Optional[str]:
    if path and os.path.isfile(path) and path.lower().endswith(".pdf"):
        return path
    # find first .pdf in CWD
    for fname in sorted(os.listdir(".")):
        if fname.lower().endswith(".pdf") and os.path.isfile(fname):
            return fname
    return None


def read_pdf_text(pdf_path: str) -> Tuple[str, List[str]]:
    """Return full_text and per-page texts (simple, no OCR)."""
    reader = PdfReader(pdf_path)
    pages = []
    for page in reader.pages:
        try:
            pages.append(page.extract_text() or "")
        except Exception:
            pages.append("")
    full_text = "\n".join(pages)
    return full_text, pages

def get_local_ollama_models(base_url: str = "http://localhost:11434") -> List[str]:
    """Get the list of locally available Ollama models."""
    try:
        r = requests.get(f"{base_url}/api/tags")
        r.raise_for_status()
        return [m["name"] for m in r.json().get("models", [])]
    except Exception:
        return []

def generate_outfile_name_with_ollama(topic: str, model: str = "llama3.1:8b") -> str:
    """Ask Ollama to propose a short filename for the PPT deck."""
    ensure_ollama_running()
    llm = ChatOllama(model=model, temperature=0.3)
    prompt = ChatPromptTemplate.from_messages([
        ("system", "You generate concise, file-safe names for PowerPoint decks."),
        ("user", f"Topic: {topic}\n\nReturn only a short file name (2–5 words) with underscores, no spaces or punctuation.")
    ])
    try:
        name = (prompt | llm).invoke({}).content.strip()
        # sanitize
        name = re.sub(r'[^A-Za-z0-9_]+', '', name)
        if not name:
            name = "presentation"
    except Exception:
        name = "presentation"
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"{name}_{timestamp}.pptx"

# =========================
# LLM: PDF → notes (map-reduce)
# =========================
def summarize_chunk_notes(llm: ChatOllama, chunk: str, max_bullets: int = 7) -> List[str]:
    prompt = ChatPromptTemplate.from_messages(
        [
            ("system",
             "You extract concise bullet notes from an academic paper section. "
             "Output ONLY a JSON array of strings. Each string <= 18 words. No nested objects."),
            ("user",
             "Section:\n---\n{chunk}\n---\n"
             f"Return a JSON array with up to {max_bullets} bullets capturing key claims, methods, results, or implications.")
        ]
    )
    raw = (prompt | llm).invoke({"chunk": chunk}).content
    # Try to parse array
    start = raw.find("["); end = raw.rfind("]")
    if start != -1 and end != -1 and end > start:
        try:
            arr = json.loads(raw[start:end+1])
            return [str(x).strip() for x in arr if str(x).strip()]
        except Exception:
            pass
    # Fallback: split lines
    lines = [ln.strip("-•* \t") for ln in raw.splitlines() if ln.strip()]
    return lines[:max_bullets]


def reduce_notes(llm: ChatOllama, bullets: List[str], max_out: int = 40) -> List[str]:
    """Compress many bullets into a focused list to feed slide planner."""
    joined = "\n".join(f"- {b}" for b in bullets)
    prompt = ChatPromptTemplate.from_messages(
        [
            ("system",
             "You condense notes from a research paper. Output ONLY a JSON array of strings. Each <= 18 words."),
            ("user",
             "Notes:\n{notes}\n\n"
             f"Return at most {max_out} distilled bullets preserving methods, datasets, key results, limitations, and takeaways.")
        ]
    )
    raw = (prompt | llm).invoke({"notes": joined}).content
    start = raw.find("["); end = raw.rfind("]")
    if start != -1 and end != -1 and end > start:
        try:
            arr = json.loads(raw[start:end+1])
            return [str(x).strip() for x in arr if str(x).strip()]
        except Exception:
            pass
    # fallback: truncate original
    return bullets[:max_out]


def extract_reference_lines(pages: List[str], tail_pages: int = 3, max_refs: int = 6) -> List[str]:
    """Heuristically pull likely references from last pages."""
    tail = "\n".join(pages[-tail_pages:])
    lines = [ln.strip() for ln in tail.splitlines() if ln.strip()]
    # crude filter: lines with year patterns or DOI/arXiv
    ref_like = []
    for ln in lines:
        if re.search(r"\b(19|20)\d{2}\b", ln) or "doi" in ln.lower() or "arxiv" in ln.lower():
            ref_like.append(ln)
    # de-duplicate & trim
    seen = set(); out = []
    for r in ref_like:
        if r not in seen:
            seen.add(r)
            out.append(r)
        if len(out) >= max_refs:
            break
    return out


# =========================
# LLM: Deck planning (paper-aware)
# =========================
def build_deck_with_llm(
    topic: str,
    audience: str = "general",
    tone: str = "informative",
    slide_count: int = 8,
    model: str = "llama3.1:8b",
    temperature: float = 0.2,
    paper_summary_bullets: Optional[List[str]] = None,
    add_references: bool = True,
    reference_lines: Optional[List[str]] = None,
) -> Deck:
    ensure_ollama_running()
    llm = ChatOllama(model=model, temperature=temperature, format="json")

    parser = PydanticOutputParser(pydantic_object=Deck)

    context_block = ""
    if paper_summary_bullets:
        context_block = "Paper context bullets:\n" + "\n".join(f"- {b}" for b in paper_summary_bullets[:80])

    prompt = ChatPromptTemplate.from_messages(
        [
            ("system",
             "You are a presentation architect. Output ONLY a JSON object for the Deck schema.\n"
             "HARD RULES:\n"
             "- EXACTLY {slide_count} slides (content-only; welcome/thanks will be added separately).\n"
             "- Each slide has: title, subtitle (<= 12 words), details (3–6 strings), optional speaker_notes.\n"
             "- Titles must be informative and unique.\n"
             "- Details must be strings (no nested objects). Each <= 14 words.\n"
             "- Ground slides in the provided paper context if given. Avoid generic textbook filler.\n"
             "- Story arc: hook, background, method/approach, data/experiments, results, limitations, implications, next steps.\n"
             "- Explain the core insight: models exhibit consistent stylistic biases exploitable by inversion models.\n"
             "- Summarize how inversion differs from embedding inversion and why it’s harder (discrete vs continuous space).\n"
             "- No prose outside the JSON."),
            ("user",
             "Topic: {topic}\nAudience: {audience}\nTone: {tone}\nSlides (no welcome/thanks): {slide_count}\n\n"
             "{context}\n\n{format_instructions}\n"
             "Additionally:\n"
             "- Slide 1 should hook with concrete problem + contribution.\n"
             "- Include one Methods/Architecture slide and one Results/Findings slide.\n"
             "- Include one Limitations/Threats to Validity slide if applicable.\n"
             "- End with clear Next Steps.")
        ]
    ).partial(format_instructions=parser.get_format_instructions(), context=context_block)

    # Get & parse plan
    raw = (prompt | llm).invoke(
        {"topic": topic, "audience": audience, "tone": tone, "slide_count": slide_count}
    ).content

    raw_json = extract_json_block(raw)
    try:
        data = json.loads(raw_json)
    except Exception:
        cleaned = re.sub(r",\s*([}\]])", r"\1", raw_json)
        data = json.loads(cleaned)

    data = repair_deck_obj(data)
    deck = Deck.model_validate(data)

    # Enforce slide_count
    if len(deck.slides) > slide_count:
        deck.slides = deck.slides[:slide_count]
    elif len(deck.slides) < slide_count:
        while len(deck.slides) < slide_count:
            deck.slides.append(
                Slide(title="Additional Insight", subtitle="Overview", details=["Key point", "Example", "Action"])
            )

    # Inject Welcome & Thank You
    welcome = Slide(
        title=f"Welcome to {deck.topic.title()}",
        subtitle=f"Exploring {deck.topic.lower()}",
        details=["Presenter: (name or team)", f"Topic: {deck.topic}", "Overview & goals"],
        speaker_notes="Greet audience; explain purpose."
    )
    thankyou = Slide(
        title="Thank You!",
        subtitle="Questions & discussion",
        details=["Contact: (email/URL)", "Resources available on request", "Appreciate your time!"],
        speaker_notes="Invite Q&A."
    )

    # Optional References (before Thank You)
    if add_references and reference_lines:
        # collapse long refs to ~1 line each
        refs = []
        for r in reference_lines:
            r = re.sub(r"\s+", " ", r)
            if len(r) > 110:
                r = r[:107] + "..."
            refs.append(r)
        refs = refs[:6] if refs else ["References available upon request"]
        references_slide = Slide(
            title="References",
            subtitle="Key citations from the paper",
            details=refs
        )
        deck.slides = [welcome] + deck.slides + [references_slide, thankyou]
    else:
        deck.slides = [welcome] + deck.slides + [thankyou]

    deck.slide_count = len(deck.slides)
    return deck


# =========================
# PowerPoint Builder (styled)
# =========================
def build_pptx(deck: Deck, outfile: str = "deck.pptx", template: Optional[str] = None):
    prs = Presentation(template) if template else Presentation()
    slide_w, slide_h = prs.slide_width, prs.slide_height

    for i, s in enumerate(deck.slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        # Special background: Welcome (0) & Thank You (last) darker
        is_hero = (i == 0) or (i == len(deck.slides) - 1)
        fill = slide.background.fill
        fill.solid()
        if is_hero:
            fill.fore_color.rgb = RGBColor(25, 25, 112)   # deep navy
            TITLE_COLOR = RGBColor(255, 255, 255)
            SUBTITLE_COLOR = RGBColor(230, 230, 230)
            BULLET_COLOR = RGBColor(240, 240, 240)
        else:
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            TITLE_COLOR = RGBColor(25, 25, 112)
            SUBTITLE_COLOR = RGBColor(70, 70, 70)
            BULLET_COLOR = RGBColor(40, 40, 40)

        # Title
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.75), slide_w - Inches(2.0), Inches(1.0))
        title_tf = title_box.text_frame
        p = title_tf.paragraphs[0]
        p.text = s.title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        title_box.width += Inches(0.3)

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.7), slide_w - Inches(2.0), Inches(0.8))
        sub_tf = sub_box.text_frame
        sp = sub_tf.paragraphs[0]
        sp.text = s.subtitle
        sp.font.size = Pt(22)
        sp.font.color.rgb = SUBTITLE_COLOR
        sp.alignment = PP_ALIGN.CENTER
        sub_box.width += Inches(0.3)

        # Details
        details_box = slide.shapes.add_textbox(Inches(1.3), Inches(2.8), slide_w - Inches(2.6), slide_h - Inches(3.3))
        tf = details_box.text_frame
        for j, d in enumerate(s.details):
            para = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
            para.text = "• " + d
            para.font.size = Pt(20)
            para.font.color.rgb = BULLET_COLOR
            para.line_spacing = 1.3
            para.space_after = Pt(6)
            para.alignment = PP_ALIGN.LEFT
        details_box.width += Inches(0.3)

        # Footer
        footer = slide.shapes.add_textbox(Inches(0.5), slide_h - Inches(0.6), slide_w - Inches(1.0), Inches(0.3))
        ftf = footer.text_frame
        fp = ftf.paragraphs[0]
        fp.text = deck.topic
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(150, 150, 150)
        fp.alignment = PP_ALIGN.RIGHT

    prs.save(outfile)
    return outfile


# =========================
# CLI
# =========================
def main():
    import argparse

    local_models = get_local_ollama_models()
    if not local_models:
        print("❌ Could not find any local Ollama models. Please run `ollama pull <model_name>`.")
        sys.exit(1)

    parser = argparse.ArgumentParser(description="Build a PowerPoint using LangChain + Ollama. Can read a research PDF.")
    parser.add_argument("--topic", required=True, help="Topic of the presentation (used in slides & footer)")
    parser.add_argument("--audience", default="general")
    parser.add_argument("--tone", default="informative")
    parser.add_argument("--slides", type=int, default=8, help="Number of content slides (welcome/thanks added automatically)")
    parser.add_argument("--model", default=None, help=f"Model to use. If not specified, defaults to the first available model. Options: {', '.join(local_models)}") # No default here, we'll set it intelligently
    parser.add_argument("--outfile", default="deck.pptx")
    parser.add_argument("--template", default=None, help="Optional .potx/.pptx theme")
    parser.add_argument("--pdf", default=None, help="Path to a research paper PDF; if omitted, auto-detects first *.pdf in CWD")
    parser.add_argument("--no-references", action="store_true", help="Disable references slide")
    args = parser.parse_args()

    if args.model:
        if args.model not in local_models:
            print(f"❌ Model '{args.model}' not found locally. Available models: {', '.join(local_models)}")
            sys.exit(1)
        selected_model = args.model
    else:
        selected_model = local_models[0]
        print(f"ℹ️ No model specified. Defaulting to '{selected_model}'.")

    paper_bullets = None
    reference_lines = None

    # Detect PDF and summarize
    pdf_path = autodetect_pdf(args.pdf)
    if pdf_path:
        print(f"📄 Using PDF: {pdf_path}")
        full_text, pages = read_pdf_text(pdf_path)
        if not full_text.strip():
            print("⚠️ PDF text appears empty (scanned/OCR?). Slides will be generic.")
        else:
            # Split into chunks ~ 1800 chars with overlap
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=1800, chunk_overlap=200, length_function=len, separators=["\n\n", "\n", " ", ""]
            )
            chunks = splitter.split_text(full_text)
            # Map: extract notes per chunk
            llm_notes = ChatOllama(model=selected_model, temperature=0.1, format="json")
            all_bullets: List[str] = []
            for ch in chunks[:24]:  # cap for speed
                notes = summarize_chunk_notes(llm_notes, ch, max_bullets=7)
                all_bullets.extend(notes)
            # Reduce: compress
            paper_bullets = reduce_notes(llm_notes, all_bullets, max_out=40)
            # References
            reference_lines = extract_reference_lines(pages, tail_pages=3, max_refs=6)
    else:
        print("ℹ️ No PDF found in CWD and --pdf not provided. Building a generic deck from topic.")

    try:
        deck = build_deck_with_llm(
            topic=args.topic,
            audience=args.audience,
            tone=args.tone,
            slide_count=args.slides,
            model=selected_model,
            paper_summary_bullets=paper_bullets,
            add_references=not args.no_references,
            reference_lines=reference_lines,
        )
    except ValidationError as ve:
        print("❌ Model returned invalid JSON for the Deck schema.")
        print(ve)
        sys.exit(2)
    except Exception as e:
        print(e)
        sys.exit(2)

    try:
        if not args.outfile or args.outfile.strip() == "deck.pptx":
            args.outfile = generate_outfile_name_with_ollama(args.topic, model=selected_model)
            print(f"🪶 Ollama named your file: {args.outfile}")

        path = build_pptx(deck, args.outfile, template=args.template)

        print(f"✅ PowerPoint created: {path}")
    except Exception as e:
        print("❌ Failed to write PowerPoint:", e)
        sys.exit(3)


if __name__ == "__main__":
    main()
